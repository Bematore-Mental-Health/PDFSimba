from flask import Flask, request, jsonify, send_from_directory, render_template, url_for, redirect
from flask_cors import CORS
import os
import threading
import time
from bs4 import BeautifulSoup
from docx import Document
from html2docx import html2docx
import uuid
import pandas as pd
from flask import Flask, request, jsonify, send_from_directory, render_template_string
from PyPDF2 import PdfReader,  PdfWriter
from flask_cors import CORS
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from flask import render_template
from urllib.parse import unquote
from werkzeug.utils import secure_filename
import tempfile
import subprocess
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from PIL import Image
import pdf2image
from docx2pdf import convert
import tempfile
from win32com import client as wc
from pdf2docx import Converter
from win32com.client import Dispatch
from werkzeug.utils import secure_filename
from bs4 import BeautifulSoup
import traceback
import pythoncom
import base64
import re
from flask import Flask, request, jsonify, render_template, send_from_directory
from flask_cors import CORS
from pdf2docx import Converter
from win32com.client import Dispatch
from werkzeug.utils import secure_filename
import time
import threading
import shutil
from docx.shared import RGBColor
from bs4 import  Tag
from docx.document import Document as DocxDocument
from docx.table import _Cell
from docx.text.paragraph import Paragraph
from docx.text.run import Run
from docx.enum.style import WD_STYLE_TYPE
from docx.shared import Pt, Inches
from docx import Document
from docx.document import Document as DocxDocument




from flask import Flask, request, jsonify, send_from_directory, send_file, render_template
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from flask import Flask, request, jsonify, send_file, make_response
from PyPDF2 import PdfMerger
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from flask import Flask, request, jsonify, g
import sqlite3
import requests 
import pytesseract
import ocrmypdf
import pdfplumber
from docx import Document
from docx.shared import Inches
import io
from flask import Flask, request, jsonify, send_from_directory
from htmldocx import HtmlToDocx




from word_to_pdf import convert_word_file
from excel_to_pdf import convert_excel_file
from powerpoint_to_pdf import convert_ppt_file
from jpg_to_pdf import convert_jpg_file
from pdf_to_word import convert_pdf_file
from pdf_to_excel import convert_pdf_to_excel
from pdf_to_ppt import convert_pdf_to_ppt 
from pdf_to_jpg import convert_pdf_to_jpg
from pdf_to_png import convert_pdf_to_png
from pdf_to_pdfa import convert_to_pdfa
from merge_pdfs import merge_pdfs
from split_pdfs import  split_pdf_logic




app = Flask(__name__)
application = app
CORS(app)


BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_FOLDER = os.path.abspath(os.path.join(BASE_DIR, '..', 'uploads'))
OUTPUT_FOLDER = os.path.abspath(os.path.join(BASE_DIR, '..', 'converted'))

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['OUTPUT_FOLDER'] = OUTPUT_FOLDER
app.config['ALLOWED_IWORK_EXTENSIONS'] = {'pages', 'key', 'numbers'}
app.config['MAX_CONTENT_LENGTH'] = None




# Word to PDF
DATABASE = os.path.join(BASE_DIR, 'conversions.db')

# Database setup
def get_db():
    db = getattr(g, '_database', None)
    if db is None:
        db = g._database = sqlite3.connect(DATABASE)
        db.row_factory = sqlite3.Row
    return db

@app.teardown_appcontext
def close_connection(exception):
    db = getattr(g, '_database', None)
    if db is not None:
        db.close()

def init_db():
    with app.app_context():
        db = get_db()
        cursor = db.cursor()
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS conversions (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id TEXT NOT NULL,
                conversion_type TEXT NOT NULL,
                original_filename TEXT NOT NULL,
                converted_filename TEXT,
                file_size INTEGER,
                status TEXT NOT NULL,
                error_message TEXT,
                conversion_id TEXT,
                parent_conversion_id TEXT,
                timestamp DATETIME DEFAULT CURRENT_TIMESTAMP
            )
        """)
        db.commit()

# Initialize database and folders
with app.app_context():
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    os.makedirs(OUTPUT_FOLDER, exist_ok=True)
    init_db()

@app.route('/convert-word-to-pdf', methods=['POST'])
def convert_word_to_pdf():
    try:
        file = request.files.get('word_file')
        user_id = request.form.get('user_id')
        
        if not file:
            return jsonify({'error': 'No file uploaded'}), 400

        ext = os.path.splitext(file.filename)[1].lower()
        if ext not in ['.doc', '.docx']:
            return jsonify({'error': 'Unsupported file format'}), 400

        conversion_id = str(uuid.uuid4())
        filename = f"{uuid.uuid4()}{ext}"
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        file.save(filepath)

        pdf_filename = f"{os.path.splitext(filename)[0]}.pdf"
        pdf_filepath = os.path.join(OUTPUT_FOLDER, pdf_filename)

        try:
            comtypes.CoInitialize()
            word = comtypes.client.CreateObject('Word.Application')
            word.Visible = False
            doc = word.Documents.Open(filepath)
            doc.SaveAs(pdf_filepath, FileFormat=17)
            doc.Close()
            word.Quit()
        except Exception as e:
            return jsonify({
                'error': str(e),
                'conversion_id': conversion_id
            }), 500
        finally:
            comtypes.CoUninitialize()

        return jsonify({
            'pdf_path': f"converted/{pdf_filename}",
            'conversion_id': conversion_id,
            'original_filename': file.filename,
            'converted_filename': pdf_filename
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/record-conversion', methods=['POST'])
def record_conversion():
    try:
        data = request.get_json()
        required_fields = ['user_id', 'conversion_type', 'original_filename', 'status']
        
        for field in required_fields:
            if field not in data:
                return jsonify({'error': f'Missing required field: {field}'}), 400
        
        db = get_db()
        cursor = db.cursor()
        cursor.execute("""
            INSERT INTO conversions 
            (user_id, conversion_type, original_filename, converted_filename, file_size, status, error_message, conversion_id) 
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        """, (
            data['user_id'],
            data['conversion_type'],
            data['original_filename'],
            data.get('converted_filename'),
            data.get('file_size'),
            data['status'],
            data.get('error_message'),
            data.get('conversion_id')
        ))
        
        db.commit()
        return jsonify({'success': True, 'id': cursor.lastrowid})
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500
        
# Excel to PDF
@app.route('/convert-excel-to-pdf', methods=['POST'])
def convert_excel_to_pdf():
    try:
        file = request.files.get('excel_file')
        user_id = request.form.get('user_id')
        
        if not file:
            return jsonify({'error': 'No file uploaded'}), 400

        ext = os.path.splitext(file.filename)[1].lower()
        if ext not in ['.xls', '.xlsx']:
            return jsonify({'error': 'Unsupported file format'}), 400

        conversion_id = str(uuid.uuid4())
        filename = f"{uuid.uuid4()}{ext}"
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        file.save(filepath)

        pdf_filename = f"{os.path.splitext(filename)[0]}.pdf"
        pdf_filepath = os.path.join(OUTPUT_FOLDER, pdf_filename)

        try:
            comtypes.CoInitialize()
            excel = comtypes.client.CreateObject('Excel.Application')
            excel.Visible = False
            workbook = excel.Workbooks.Open(filepath)
            
            # Format settings for better PDF output
            for sheet in workbook.Sheets:
                sheet.PageSetup.Orientation = 2  # Landscape
                sheet.PageSetup.Zoom = False
                sheet.PageSetup.FitToPagesWide = 1
                sheet.PageSetup.FitToPagesTall = False
                sheet.PageSetup.CenterHorizontally = True
                sheet.PageSetup.CenterVertically = True
            
            workbook.ExportAsFixedFormat(0, pdf_filepath)
            workbook.Close(False)
            excel.Quit()
        except Exception as e:
            return jsonify({
                'error': str(e),
                'conversion_id': conversion_id
            }), 500
        finally:
            comtypes.CoUninitialize()

        return jsonify({
            'pdf_path': f"converted/{pdf_filename}",
            'conversion_id': conversion_id,
            'original_filename': file.filename,
            'converted_filename': pdf_filename
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500


# PowerPoint to PDF
@app.route('/convert-ppt-to-pdf', methods=['POST'])
def convert_ppt_to_pdf():
    try:
        file = request.files.get('ppt_file')
        user_id = request.form.get('user_id')
        
        if not file:
            return jsonify({'error': 'No file uploaded'}), 400

        ext = os.path.splitext(file.filename)[1].lower()
        if ext not in ['.ppt', '.pptx']:
            return jsonify({'error': 'Unsupported file format'}), 400

        conversion_id = str(uuid.uuid4())
        filename = f"{uuid.uuid4()}{ext}"
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        file.save(filepath)

        pdf_filename = f"{os.path.splitext(filename)[0]}.pdf"
        pdf_filepath = os.path.join(OUTPUT_FOLDER, pdf_filename)

        try:
            comtypes.CoInitialize()
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1  # Set to 1 for visible, 0 for invisible
            presentation = powerpoint.Presentations.Open(filepath, WithWindow=False)
            presentation.SaveAs(pdf_filepath, 32)  # 32 is the PDF format
            presentation.Close()
            powerpoint.Quit()
        except Exception as e:
            return jsonify({
                'error': str(e),
                'conversion_id': conversion_id
            }), 500
        finally:
            comtypes.CoUninitialize()

        return jsonify({
            'pdf_path': f"converted/{pdf_filename}",
            'conversion_id': conversion_id,
            'original_filename': file.filename,
            'converted_filename': pdf_filename
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500

# JPG/PNG to PDF
@app.route('/convert-jpg-to-pdf', methods=['POST'])
def convert_jpg_to_pdf():
    try:
        files = request.files.getlist('jpg_file')
        user_id = request.form.get('user_id')
        
        if not files or len(files) == 0:
            return jsonify({'error': 'No files uploaded'}), 400

        conversion_id = str(uuid.uuid4())
        temp_files = []
        images = []
        total_size = 0

        try:
            for file in files:
                ext = os.path.splitext(file.filename)[1].lower()
                if ext not in ['.jpg', '.jpeg', '.png']:
                    return jsonify({'error': 'Unsupported file format'}), 400

                filename = f"{uuid.uuid4()}{ext}"
                filepath = os.path.join(UPLOAD_FOLDER, filename)
                file.save(filepath)
                temp_files.append(filepath)
                total_size += os.path.getsize(filepath)

                img = Image.open(filepath).convert("RGB")
                images.append(img)

            if not images:
                return jsonify({'error': 'No valid images found'}), 400

            pdf_filename = f"{uuid.uuid4()}.pdf"
            pdf_filepath = os.path.join(OUTPUT_FOLDER, pdf_filename)

            # Save all images to PDF
            images[0].save(
                pdf_filepath,
                save_all=True,
                append_images=images[1:],
                quality=100,
                optimize=True
            )

            # Return response (database recording will be done by frontend)
            return jsonify({
                'pdf_path': f"converted/{pdf_filename}",
                'conversion_id': conversion_id,
                'original_filename': ', '.join([f.filename for f in files]),
                'converted_filename': pdf_filename
            })

        except Exception as e:
            return jsonify({
                'error': str(e),
                'conversion_id': conversion_id
            }), 500
        finally:
            # Clean up temporary files
            for temp_file in temp_files:
                try:
                    os.remove(temp_file)
                except:
                    pass

    except Exception as e:
        return jsonify({'error': str(e)}), 500

# PDF To Word
@app.route('/convert-pdf-to-word', methods=['POST'])
def convert_pdf_to_word_route():
    try:
        file = request.files.get('pdf_file')
        user_id = request.form.get('user_id')
        
        if not file:
            return jsonify({'error': 'No file uploaded'}), 400

        ext = os.path.splitext(file.filename)[1].lower()
        if ext != '.pdf':
            return jsonify({'error': 'Only PDF files are supported'}), 400

        conversion_id = str(uuid.uuid4())
        filename = f"{uuid.uuid4()}.pdf"
        input_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(input_path)

        base_name = os.path.splitext(filename)[0]
        html_filename = f"{base_name}.html"
        editor_html_filename = f"editor_{base_name}.html"
        docx_filename = f"{base_name}.docx"
        html_path = os.path.join(OUTPUT_FOLDER, html_filename)
        editor_html_path = os.path.join(OUTPUT_FOLDER, editor_html_filename)
        docx_path = os.path.join(OUTPUT_FOLDER, docx_filename)

        try:
            # Create a Word document
            doc = Document()
            
            # Convert PDF to images (one per page)
            images = convert_from_path(input_path)
            
            html_content = "<html><body>"
            editor_html_content = "<html><body>"
            
            for i, image in enumerate(images):
                # Save image temporarily
                img_path = os.path.join(UPLOAD_FOLDER, f"temp_{i}.jpg")
                image.save(img_path, 'JPEG')
                
                # Extract text from image using OCR
                text = pytesseract.image_to_string(Image.open(img_path))
                
                # Add text to Word document
                if text.strip():
                    doc.add_paragraph(text)
                
                # Add image to Word document
                doc.add_picture(img_path, width=Inches(6))
                
                # Add page break if not last page
                if i < len(images) - 1:
                    doc.add_page_break()
                
                # Build HTML content for preview (with images)
                html_content += f"<div><p>{text.replace('\n', '<br>')}</p>"
                html_content += f"<img src='data:image/jpeg;base64,{image_to_base64(image)}' style='max-width: 100%;'/></div>"
                
                # Build HTML for editor (without images)
                editor_html_content += f"<div><p>{text.replace('\n', '<br>')}</p>"
                editor_html_content += f"<div class='image-placeholder' style='border: 1px dashed #ccc; padding: 10px; margin: 10px 0;'>[Image - preserved in final document]</div></div>"
                
                # Clean up temp image
                os.remove(img_path)
            
            html_content += "</body></html>"
            editor_html_content += "</body></html>"
            
            # Save both HTML versions
            with open(html_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            with open(editor_html_path, 'w', encoding='utf-8') as f:
                f.write(editor_html_content)
            
            # Save Word document
            doc.save(docx_path)
            
            return jsonify({
                'word_path': f"converted/{html_filename}",
                'word_docx_path': f"converted/{docx_filename}",
                'word_filename': html_filename,
                'conversion_id': conversion_id,
                'original_filename': file.filename,
                'converted_filename': docx_filename
            })

        except Exception as e:
            return jsonify({
                'error': f'PDF conversion failed: {str(e)}',
                'conversion_id': conversion_id
            }), 500

    except Exception as e:
        return jsonify({'error': str(e)}), 500

def image_to_base64(image):
    """Convert PIL image to base64 string"""
    buffered = io.BytesIO()
    image.save(buffered, format="JPEG")
    return base64.b64encode(buffered.getvalue()).decode('utf-8')

@app.route('/editor/<filename>')
def edit_word_document(filename):
    return f"""
    <!DOCTYPE html>
    <html>
    <head>
      <title>PDFSimba| Edit Word Document</title>
      <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.0/dist/css/bootstrap.min.css" rel="stylesheet">
      <link href="https://cdn.jsdelivr.net/npm/summernote@0.8.18/dist/summernote-lite.min.css" rel="stylesheet">
      <link rel="icon" href="../Media/book3.png" type="image/x-icon">
      <style>
        body {{
          background-color: #f8f9fa;
          padding: 30px;
        }}
        .editor-container {{
          background: #fff;
          padding: 20px;
          border-radius: 10px;
          box-shadow: 0 0 10px rgba(0,0,0,0.1);
        }}
        .image-placeholder {{
          background-color: #f8f9fa;
          color: #6c757d;
          text-align: center;
        }}
      </style>
    </head>
    <body>
      <div class="container">
        <div class="editor-container">
          <h3 class="mb-4">Edit Your Word Document</h3>
          
          <form method="POST" action="/save-edited-word/{filename}">
            <button class="btn btn-success mb-3" type="submit">Save and Download</button> 
            <textarea id="editor" name="content"></textarea>
          </form>
        </div>
      </div>

      <script src="https://code.jquery.com/jquery-3.6.0.min.js"></script>
      <script src="https://cdn.jsdelivr.net/npm/bootstrap@5.3.0/dist/js/bootstrap.bundle.min.js"></script>
      <script src="https://cdn.jsdelivr.net/npm/summernote@0.8.18/dist/summernote-lite.min.js"></script>

      <script>
        $(document).ready(function() {{
          fetch("/converted/editor_{filename}")
            .then(response => {{
              if (!response.ok) throw new Error('Failed to load content');
              return response.text();
            }})
            .then(content => {{
              $('#editor').summernote({{
                height: 500,
                tabsize: 2,
                toolbar: [
                  ['style', ['bold', 'italic', 'underline', 'clear']],
                  ['font', ['strikethrough', 'superscript', 'subscript']],
                  ['fontsize', ['fontsize']],
                  ['color', ['color']],
                  ['para', ['ul', 'ol', 'paragraph']],
                  ['insert', ['link']],  // No image upload option
                  ['view', ['fullscreen', 'codeview', 'help']]
                ]
              }});
              $('#editor').summernote('code', content);
            }})
            .catch(error => {{
              console.error('Error:', error);
              $('#editor').summernote({{
                height: 500,
                tabsize: 2
              }});
              $('#editor').summernote('code', '<p>Document loaded without images for better performance</p>');
            }});
        }});
      </script>
    </body>
    </html>
    """

@app.route('/save-edited-word/<filename>', methods=['POST'])
def save_edited_word(filename):
    content = request.form.get('content')
    if not content:
        return "No content provided", 400

    try:
        # Load the original HTML with images
        original_html_path = os.path.join(OUTPUT_FOLDER, filename)
        if not os.path.exists(original_html_path):
            return "Original document not found", 404

        with open(original_html_path, 'r', encoding='utf-8') as f:
            original_html = f.read()

        # Parse both HTML contents
        edited_soup = BeautifulSoup(content, 'html.parser')
        original_soup = BeautifulSoup(original_html, 'html.parser')

        # Replace placeholders with original images
        edited_placeholders = edited_soup.select('.image-placeholder')
        original_images = original_soup.select('img')

        if len(edited_placeholders) != len(original_images):
            return "Mismatch between edited content and original images", 400

        for placeholder, img in zip(edited_placeholders, original_images):
            placeholder.replace_with(img)

        # Create the final document
        doc = Document()
        new_parser = HtmlToDocx()
        new_parser.add_html_to_document(str(edited_soup), doc)

        # Save the document
        output_filename = f"edited_{os.path.splitext(filename)[0]}.docx"
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        doc.save(output_path)

        return send_from_directory(
            OUTPUT_FOLDER,
            output_filename,
            as_attachment=True,
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )

    except Exception as e:
        return f"Error saving Word file: {str(e)}", 500
        
# PDF to Excel
import os
import uuid
import bisect
from collections import Counter
from flask import request, jsonify
import pandas as pd
from PyPDF2 import PdfReader
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
try:
    import camelot
    CAMELOT_AVAILABLE = True
except ImportError:
    CAMELOT_AVAILABLE = False

@app.route('/convert-pdf-to-excel', methods=['POST'])
def convert_pdf_to_excel():
    file = request.files.get('pdf_file')
    user_id = request.form.get('user_id')
    conversion_id = str(uuid.uuid4())
    
    if not file:
        return jsonify({
            'error': 'No file uploaded',
            'conversion_id': conversion_id
        }), 400

    # Create necessary directories if they don't exist
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    os.makedirs(OUTPUT_FOLDER, exist_ok=True)

    filename = f"{uuid.uuid4()}.pdf"
    input_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(input_path)

    output_filename = f"{os.path.splitext(filename)[0]}.xlsx"
    output_path = os.path.join(OUTPUT_FOLDER, output_filename)

    try:
        # First try direct text extraction
        try:
            with open(input_path, 'rb') as f:
                reader = PdfReader(f)
                text = ''.join(page.extract_text() or "" for page in reader.pages)
                
                if text.strip():
                    # Try to detect tabular structure
                    if '\t' in text or any(len(line.split()) > 3 for line in text.splitlines()):
                        rows = [line.split('\t') if '\t' in line else line.split() 
                               for line in text.splitlines() if line.strip()]
                    else:
                        # Plain text - each paragraph becomes a row
                        rows = [[paragraph] for paragraph in text.split('\n\n') if paragraph.strip()]
                    
                    pd.DataFrame(rows).to_excel(output_path, index=False, header=False)
                    return jsonify({
                        'excel_filename': output_filename,
                        'conversion_id': conversion_id,
                        'original_filename': file.filename
                    })
        except Exception as e:
            print(f"Direct extraction failed: {str(e)}")

        # If Camelot is available, try table extraction
        if CAMELOT_AVAILABLE:
            try:
                tables = camelot.read_pdf(input_path, flavor='lattice', pages='all')
                if tables and len(tables) > 0:
                    df = pd.concat([t.df for t in tables])
                    df.to_excel(output_path, index=False, header=False)
                    return jsonify({
                        'excel_filename': output_filename,
                        'conversion_id': conversion_id,
                        'original_filename': file.filename
                    })
            except Exception as e:
                print(f"Camelot extraction failed: {str(e)}")

        # Fall back to OCR
        try:
            images = convert_from_path(input_path)
            all_rows = []
            
            for i, image in enumerate(images):
                img_path = os.path.join(UPLOAD_FOLDER, f"temp_{i}.jpg")
                image.save(img_path, 'JPEG')
                
                # Try table detection
                try:
                    ocr_data = pytesseract.image_to_data(
                        Image.open(img_path),
                        config='--psm 6',
                        output_type=pytesseract.Output.DICT
                    )
                    rows = process_ocr_to_table(ocr_data)
                    if len(rows) > 1:  # If we found a table structure
                        all_rows.extend(rows)
                        continue
                except Exception as e:
                    print(f"OCR table detection failed: {str(e)}")
                
                # Fall back to line-by-line extraction
                text = pytesseract.image_to_string(Image.open(img_path))
                if text.strip():
                    all_rows.extend([[line] for line in text.splitlines() if line.strip()])
                
                os.remove(img_path)
            
            pd.DataFrame(all_rows).to_excel(output_path, index=False, header=False)
            return jsonify({
                'excel_filename': output_filename,
                'conversion_id': conversion_id,
                'original_filename': file.filename
            })
        except Exception as e:
            print(f"OCR extraction failed: {str(e)}")
            raise

    except Exception as e:
        return jsonify({
            'error': f'PDF to Excel conversion failed: {str(e)}',
            'conversion_id': conversion_id
        }), 500
    finally:
        if os.path.exists(input_path):
            os.remove(input_path)

def process_ocr_to_table(ocr_data):
    """Convert OCR output to table structure preserving alignment"""
    # Group by line and word position
    lines = {}
    for i in range(len(ocr_data['text'])):
        if ocr_data['text'][i].strip():
            line_num = ocr_data['line_num'][i]
            left_pos = ocr_data['left'][i]
            if line_num not in lines:
                lines[line_num] = []
            lines[line_num].append((left_pos, ocr_data['text'][i]))
    
    if not lines:
        return []
    
    # Find common left positions as column dividers
    all_positions = [pos for line in lines.values() for pos, _ in line]
    position_counts = Counter(all_positions)
    common_positions = [pos for pos, cnt in position_counts.most_common(5) if cnt > 1]
    common_positions.sort()
    
    # Create table data
    table_data = []
    for line_num in sorted(lines.keys()):
        row = [''] * len(common_positions) if common_positions else []
        for pos, text in lines[line_num]:
            if common_positions:
                col = bisect.bisect_right(common_positions, pos) - 1
                if 0 <= col < len(row):
                    row[col] = text if not row[col] else row[col] + ' ' + text
            else:
                row.append(text)
        table_data.append(row)
    
    return table_data

@app.route('/download/<filename>')
def download_excel(filename):
    return send_from_directory(OUTPUT_FOLDER, filename, as_attachment=True)

@app.route('/converted/<path:filename>')
def serve_pdf(filename):
    return send_from_directory(OUTPUT_FOLDER, filename, as_attachment=True)



@app.route('/excel-editor/<filename>')
def edit_excel(filename):
    return f"""
    <!DOCTYPE html>
    <html>
    <head>
        <title>PDFSimba| Edit Excel</title>
        <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/handsontable@13.0.0/dist/handsontable.full.min.css">
        <script src="https://cdn.jsdelivr.net/npm/handsontable@13.0.0/dist/handsontable.full.min.js"></script>
        <script src="https://cdnjs.cloudflare.com/ajax/libs/xlsx/0.18.5/xlsx.full.min.js"></script>
        <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.3/dist/css/bootstrap.min.css" rel="stylesheet"/>
        <style>
            body {{ margin: 20px; font-family: sans-serif; background-color: #f9f9f9; }}
            #excelEditor {{ width: 100%; height: 500px; margin-bottom: 20px; }}
            .ht_master .wtHolder {{ height: 500px !important; }}
        </style>
    </head>
    <body>
        <div class="container">
            <h2>Edit Excel File</h2>
            <button onclick="downloadExcel()" class="btn btn-success mb-2">Save and Download</button>
            <div id="excelEditor"></div>
        </div>

        <script>
            let container = document.getElementById('excelEditor');
            let hot;

            fetch("/converted/{filename}")
                .then(res => res.arrayBuffer())
                .then(buffer => {{
                    const workbook = XLSX.read(buffer, {{ type: "array" }});
                    const ws = workbook.Sheets[workbook.SheetNames[0]];
                    const data = XLSX.utils.sheet_to_json(ws, {{ header: 1 }});
                    
                    hot = new Handsontable(container, {{
                        data: data,
                        rowHeaders: true,
                        colHeaders: true,
                        width: '100%',
                        height: 500,
                        licenseKey: 'non-commercial-and-evaluation',
                        manualColumnResize: true,
                        manualRowResize: true,
                        contextMenu: true,
                        filters: true,
                        dropdownMenu: true
                    }});
                }})
                .catch(error => {{
                    console.error('Error loading Excel:', error);
                    container.innerHTML = '<div class="alert alert-danger">Error loading Excel file</div>';
                }});

            function downloadExcel() {{
                const wb = XLSX.utils.book_new();
                const ws = XLSX.utils.aoa_to_sheet(hot.getData());
                XLSX.utils.book_append_sheet(wb, ws, "Sheet1");
                XLSX.writeFile(wb, "edited_{filename}");
            }}
        </script>
    </body>
    </html>
    """


# PDF to Powerpoint
@app.route('/convert-pdf-to-ppt', methods=['POST'])
def convert_pdf_to_ppt_route():
    file = request.files.get('pdf_file')
    user_id = request.form.get('user_id')
    conversion_id = str(uuid.uuid4())
    
    ppt_filename, error = convert_pdf_to_ppt(file)
    if error:
        return jsonify({
            'error': error,
            'conversion_id': conversion_id
        }), 400
        
    return jsonify({
        'ppt_filename': ppt_filename,
        'conversion_id': conversion_id
    })

def convert_pdf_to_ppt(file):
    if not file:
        return None, "No file uploaded."

    ext = os.path.splitext(file.filename)[1].lower()
    if ext != '.pdf':
        return None, "Only PDF files are supported."

    filename = f"{uuid.uuid4()}.pdf"
    input_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(input_path)

    output_filename = f"{os.path.splitext(filename)[0]}.pptx"
    output_path = os.path.join(OUTPUT_FOLDER, output_filename)

    try:
        # Convert PDF to images
        images = convert_from_path(input_path, dpi=300)
        
        # Create PowerPoint presentation
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        
        for image in images:
            # Create temp image file
            temp_img_path = os.path.join(UPLOAD_FOLDER, f"temp_{uuid.uuid4()}.jpg")
            image.save(temp_img_path, "JPEG", quality=95)
            
            # Add slide with image
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(
                temp_img_path, 
                Inches(0), 
                Inches(0), 
                width=prs.slide_width
            )
            
            # Clean up temp image
            os.remove(temp_img_path)

        # Save presentation
        prs.save(output_path)
        
        # Clean up input PDF
        os.remove(input_path)
        
        return output_filename, None
        
    except Exception as e:
        # Clean up any remaining files
        if os.path.exists(input_path):
            os.remove(input_path)
        return None, str(e)

@app.route('/download/<filename>')
def download_file(filename):
    return send_from_directory(OUTPUT_FOLDER, filename, as_attachment=True)

def convert_pdf_to_ppt(file):
    if not file:
        return None, "No file uploaded."

    ext = os.path.splitext(file.filename)[1].lower()
    if ext != '.pdf':
        return None, "Only PDF files are supported."

    filename = f"{uuid.uuid4()}.pdf"
    input_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(input_path)

    output_filename = f"{os.path.splitext(filename)[0]}.pptx"
    output_path = os.path.join(OUTPUT_FOLDER, output_filename)

    try:
        images = convert_from_path(input_path)
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]  

        for image in images:
            slide = prs.slides.add_slide(blank_slide_layout)
            image_path = os.path.join(UPLOAD_FOLDER, f"temp_{uuid.uuid4()}.jpg")
            image.save(image_path, "JPEG")
            slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=prs.slide_width)
            os.remove(image_path)

        prs.save(output_path)
        return output_filename, None
    except Exception as e:
        return None, str(e)


# PDF to JPG 
from flask import Flask, request, jsonify, send_from_directory
from pdf_to_jpg import convert_pdf_to_jpg

import uuid
import zipfile
from flask import jsonify, request

@app.route('/convert-pdf-to-jpg', methods=['POST'])
def convert_pdf_to_jpg_route():
    file = request.files.get('pdf_file')
    user_id = request.form.get('user_id')
    
    if not file:
        return jsonify({'error': 'No file uploaded'}), 400

    conversion_id = str(uuid.uuid4())
    pdf_filename = f"{uuid.uuid4()}.pdf"
    pdf_path = os.path.join(UPLOAD_FOLDER, pdf_filename)
    file.save(pdf_path)

    image_filenames = []
    try:
        # Convert PDF to images
        images = convert_from_path(pdf_path)
        for i, image in enumerate(images):
            image_filename = f"{uuid.uuid4()}.jpg"
            image_path = os.path.join(OUTPUT_FOLDER, image_filename)
            image.save(image_path, 'JPEG')
            image_filenames.append(image_filename)

        # Create ZIP file
        zip_filename = f"{uuid.uuid4()}.zip"
        zip_path = os.path.join(OUTPUT_FOLDER, zip_filename)
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for image_filename in image_filenames:
                image_path = os.path.join(OUTPUT_FOLDER, image_filename)
                zipf.write(image_path, os.path.basename(image_path))

        image_urls = [f"/converted/{filename}" for filename in image_filenames]
        return jsonify({
            'imageUrls': image_urls, 
            'zipUrl': f"/converted/{zip_filename}",
            'zipFilename': f"{os.path.splitext(file.filename)[0]}_converted.zip",
            'conversion_id': conversion_id
        })
    except Exception as e:
        return jsonify({
            'error': f'Conversion failed: {str(e)}',
            'conversion_id': conversion_id
        }), 500
    finally:
        # Clean up the uploaded PDF file
        if os.path.exists(pdf_path):
            try:
                os.remove(pdf_path)
            except:
                pass
@app.route('/converted/<filename>')
def serve_converted_files(filename):
    try:
        return send_from_directory(OUTPUT_FOLDER, filename)
    except FileNotFoundError:
        return jsonify({'error': 'File not found'}), 404
    except Exception as e:
        return jsonify({'error': str(e)}), 500
        

# PDF To PNG
@app.route('/convert-pdf-to-png', methods=['POST'])
def convert_pdf_to_png_route():
    file = request.files.get('pdf_file')
    user_id = request.form.get('user_id')
    
    if not file:
        return jsonify({'error': 'No file uploaded'}), 400

    conversion_id = str(uuid.uuid4())
    pdf_filename = f"{uuid.uuid4()}.pdf"
    pdf_path = os.path.join(UPLOAD_FOLDER, pdf_filename)
    file.save(pdf_path)

    image_filenames = []
    try:
        images = convert_from_path(pdf_path)
        for i, image in enumerate(images):
            image_filename = f"{uuid.uuid4()}.png"
            image_path = os.path.join(OUTPUT_FOLDER, image_filename)
            image.save(image_path, 'PNG')
            image_filenames.append(image_filename)

        zip_filename = f"{uuid.uuid4()}.zip"
        zip_path = os.path.join(OUTPUT_FOLDER, zip_filename)
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for image_filename in image_filenames:
                zipf.write(os.path.join(OUTPUT_FOLDER, image_filename), image_filename)

        image_urls = [f"/converted/{filename}" for filename in image_filenames]
        return jsonify({
            'imageUrls': image_urls, 
            'zipUrl': f"/converted/{zip_filename}",
            'zipFilename': f"{os.path.splitext(file.filename)[0]}_converted.zip",
            'conversion_id': conversion_id
        })
    except Exception as e:
        return jsonify({
            'error': f'Conversion failed: {str(e)}',
            'conversion_id': conversion_id
        }), 500
    finally:
        if os.path.exists(pdf_path):
            try:
                os.remove(pdf_path)
            except:
                pass


@app.route('/converted/<filename>')
def serve_converted_png_files(filename):
    try:
        return send_from_directory(OUTPUT_FOLDER, filename)
    except FileNotFoundError:
        return jsonify({'error': 'File not found'}), 404
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# PDF To PDF/A
def convert_to_pdfa(input_path, output_path, pdfa_version):
    """Convert PDF to PDF/A using Ghostscript"""
    try:
        # Validate PDF/A version
        valid_versions = ['1A', '1B', '2A', '2B', '2U', '3A', '3B', '3U']
        if pdfa_version not in valid_versions:
            raise ValueError(f"Invalid PDF/A version. Must be one of: {', '.join(valid_versions)}")

        # Ghostscript command for PDF/A conversion
        cmd = [
            'gswin64',
            '-dPDFA',
            f'-dPDFACompatibilityPolicy={pdfa_version}',
            '-dBATCH',
            '-dNOPAUSE',
            '-dNOOUTERSAVE',
            '-dUseCIEColor',
            '-sProcessColorModel=DeviceRGB',
            '-sDEVICE=pdfwrite',
            f'-sOutputFile={output_path}',
            input_path
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode != 0:
            raise Exception(f"Ghostscript error: {result.stderr}")
            
        return True
    except Exception as e:
        raise Exception(f"Conversion failed: {str(e)}")

@app.route('/convert-to-pdfa', methods=['POST'])
def convert_pdf_to_pdfa():
    try:
        if 'pdf_file' not in request.files:
            return jsonify({'error': 'No file uploaded', 'conversion_id': str(uuid.uuid4())}), 400
            
        file = request.files['pdf_file']
        pdfa_version = request.form.get('pdfa_version', '1B')
        
        if file.filename == '':
            return jsonify({'error': 'No selected file', 'conversion_id': str(uuid.uuid4())}), 400

        original_filename = secure_filename(file.filename)
        temp_filename = f"temp_{uuid.uuid4()}_{original_filename}"
        converted_filename = f"pdfa_{pdfa_version}_{uuid.uuid4()}_{original_filename}"
        
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], temp_filename)
        converted_path = os.path.join(app.config['OUTPUT_FOLDER'], converted_filename)
        
        file.save(temp_path)
        
        try:
            convert_to_pdfa(temp_path, converted_path, pdfa_version)
            
            if not os.path.exists(converted_path):
                raise Exception("Output file not created")
            
            return jsonify({
                'success': True,
                'filename': converted_filename,
                'pdfa_version': pdfa_version,
                'download_url': url_for('download_pdfa', filename=converted_filename, _external=True),
                'conversion_id': str(uuid.uuid4())
            })
            
        except Exception as conversion_error:
            raise Exception(f"Conversion error: {str(conversion_error)}")
            
    except Exception as e:
        return jsonify({
            'error': str(e),
            'conversion_id': str(uuid.uuid4())
        }), 500
    finally:
        if 'temp_path' in locals() and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except:
                pass

@app.route('/download-pdfa/<filename>')
def download_pdfa(filename):
    try:
        file_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
        
        if not os.path.exists(file_path):
            return jsonify({'error': 'File not found'}), 404
            
        # Force download with proper headers
        return send_file(
            file_path,
            as_attachment=True,
            download_name=filename,
            mimetype='application/pdf'
        )
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Merge PDFs
@app.route('/merge-pdfs', methods=['POST'])
def merge_pdfs_route():
    files = request.files.getlist('pdf_files')
    user_id = request.form.get('user_id')
    conversion_id = str(uuid.uuid4())
    
    if len(files) < 2:
        return jsonify({
            'error': 'At least two PDF files are required',
            'conversion_id': conversion_id
        }), 400

    try:
        os.makedirs(UPLOAD_FOLDER, exist_ok=True)
        os.makedirs(OUTPUT_FOLDER, exist_ok=True)
        
        output_path = merge_pdfs(files, UPLOAD_FOLDER, OUTPUT_FOLDER)
        filename = os.path.basename(output_path)
        
        return jsonify({
            'pdf_path': f'/output/{filename}',
            'filename': filename,
            'conversion_id': conversion_id
        })
    except Exception as e:
        app.logger.error(f"Error merging PDFs: {str(e)}")
        return jsonify({
            'error': str(e),
            'conversion_id': conversion_id
        }), 500

@app.route('/output/<filename>')
def download_merged_pdf(filename):
    try:
        file_path = os.path.join(OUTPUT_FOLDER, filename)
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File {filename} not found")
            
        if request.args.get('download'):
            return send_file(
                file_path,
                mimetype='application/pdf',
                as_attachment=True,
                download_name=filename
            )
        else:
            return send_file(
                file_path,
                mimetype='application/pdf',
                as_attachment=False
            )
    except Exception as e:
        app.logger.error(f"Error serving merged PDF: {str(e)}")
        return jsonify({'error': str(e)}), 404
        
# Split PDFs
@app.route('/upload-pdf', methods=['POST'])
def upload_pdf():
    file = request.files.get('pdf_file')
    user_id = request.form.get('user_id')
    conversion_id = str(uuid.uuid4())
    
    if not file:
        return jsonify({
            'error': 'No file uploaded',
            'conversion_id': conversion_id
        }), 400

    try:
        original_filename = secure_filename(file.filename)  
        filename = f"{uuid.uuid4()}_{original_filename}"
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        
        os.makedirs(UPLOAD_FOLDER, exist_ok=True)
        file.save(file_path)

        return jsonify({
            'file_path': filename,  
            'original_filename': original_filename,
            'conversion_id': conversion_id
        })
    except Exception as e:
        return jsonify({
            'error': f"Failed to upload file: {str(e)}",
            'conversion_id': conversion_id
        }), 500


@app.route('/split-editor', methods=['GET'])
def split_editor():
    filename = request.args.get('file')  
    if not filename:
        return jsonify({'error': 'No file specified'}), 400
    
    # URL decode the filename
    decoded_filename = unquote(filename)
    
    # Check if file exists in uploads directory
    file_path = os.path.join(UPLOAD_FOLDER, decoded_filename)
    if not os.path.exists(file_path):
        app.logger.error(f"File not found: {file_path}")
        return jsonify({'error': 'File not found'}), 404
    
    # Pass the properly encoded URL to the template
    return render_template('split_editor.html', file_url=url_for('serve_uploaded_file', filename=filename))

@app.route('/uploads/<filename>')
def serve_uploaded_file(filename):
    return send_from_directory(UPLOAD_FOLDER, filename)

@app.route('/output/<filename>')
def serve_output_file(filename):
    return send_from_directory(OUTPUT_FOLDER, filename)


@app.route('/split-pdf', methods=['POST'])
def split_pdf():
    data = request.json
    filename = data.get('filePath')
    user_id = data.get('user_id')
    conversion_id = str(uuid.uuid4())
    
    if not filename:
        return jsonify({
            'error': 'No filename provided',
            'conversion_id': conversion_id
        }), 400
    
    decoded_filename = unquote(filename)
    file_path = os.path.join(UPLOAD_FOLDER, decoded_filename)
    
    if not os.path.exists(file_path):
        return jsonify({
            'error': f"File not found at path: {file_path}",
            'conversion_id': conversion_id
        }), 400

    try:
        options = {
            'splitAfter': data.get('splitAfter'),
            'rangeStart': data.get('rangeStart'),
            'rangeEnd': data.get('rangeEnd'),
            'customPages': data.get('customPages'),
            'oddPages': data.get('oddPages'),
            'evenPages': data.get('evenPages'),
        }
        output_files = split_pdf_logic(file_path, options, OUTPUT_FOLDER)
        
        # Record successful split (fire-and-forget)
        if user_id:
            try:
                num_output_files = len(output_files)
                requests.post(f"{BACKEND_URL}/record-conversion", json={
                    'user_id': user_id,
                    'conversion_type': 'split_pdf',
                    'original_filename': decoded_filename,
                    'converted_filename': f"{num_output_files}_split_files",
                    'status': 'completed',
                    'conversion_id': conversion_id,
                    'num_files_generated': num_output_files
                }, timeout=2)
            except:
                pass

        return jsonify({
            'output_files': [url_for('serve_output_file', filename=os.path.basename(f)) for f in output_files],
            'conversion_id': conversion_id
        })
    except Exception as e:
        # Record failed split (fire-and-forget)
        if user_id:
            try:
                requests.post(f"{BACKEND_URL}/record-conversion", json={
                    'user_id': user_id,
                    'conversion_type': 'split_pdf',
                    'original_filename': decoded_filename,
                    'status': 'failed',
                    'error_message': str(e)[:255],
                    'conversion_id': conversion_id
                }, timeout=2)
            except:
                pass
        
        return jsonify({
            'error': f"Failed to process PDF: {str(e)}",
            'conversion_id': conversion_id
        }), 500


# Protect PDF 
@app.route('/protect-pdf', methods=['POST'])
def protect_pdf():
    file = request.files.get('pdf_file')
    password = request.form.get('password')
    user_id = request.form.get('user_id')

    if not file or not password:
        return jsonify({'error': 'File and password are required'}), 400

    try:
        original_filename = secure_filename(file.filename)
        temp_filename = f"{uuid.uuid4()}_{original_filename}"
        temp_path = os.path.join(UPLOAD_FOLDER, temp_filename)
        protected_filename = f"protected_{uuid.uuid4()}_{original_filename}"
        protected_path = os.path.join(OUTPUT_FOLDER, protected_filename)
        conversion_id = str(uuid.uuid4())

        file.save(temp_path)

        reader = PdfReader(temp_path)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        writer.encrypt(password)

        with open(protected_path, "wb") as output_file:
            writer.write(output_file)

        os.remove(temp_path)

        return jsonify({
            'success': True,
            'filename': protected_filename,
            'download_url': url_for('download_protected', filename=protected_filename, _external=True),
            'conversion_id': conversion_id
        })

    except Exception as e:
        return jsonify({
            'error': str(e),
            'conversion_id': conversion_id if 'conversion_id' in locals() else None
        }), 500
        
@app.route('/download-protected/<filename>')
def download_protected(filename):
    try:
        file_path = os.path.join(OUTPUT_FOLDER, filename)
        if not os.path.exists(file_path):
            return jsonify({'error': 'File not found'}), 404
        return send_from_directory(OUTPUT_FOLDER, filename, as_attachment=True)
    except Exception as e:
        return jsonify({'error': str(e)}), 404


# Unlock PDF 

@app.route('/unlock-pdf', methods=['POST'])
def unlock_pdf():
    try:
        if 'pdf_file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400

        file = request.files['pdf_file']
        password = request.form.get('password')
        user_id = request.form.get('user_id')
        conversion_id = str(uuid.uuid4())

        if not password:
            return jsonify({'error': 'Password is required', 'conversion_id': conversion_id}), 400

        if file.filename == '':
            return jsonify({'error': 'No selected file', 'conversion_id': conversion_id}), 400

        # Generate unique filenames
        original_filename = secure_filename(file.filename)
        temp_filename = f"temp_{uuid.uuid4()}_{original_filename}"
        unlocked_filename = f"unlocked_{uuid.uuid4()}_{original_filename}"

        # Define paths for temporary and output files
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], temp_filename)
        unlocked_path = os.path.join(app.config['OUTPUT_FOLDER'], unlocked_filename)

        # Save uploaded file
        file.save(temp_path)

        try:
            # Process the uploaded PDF
            reader = PdfReader(temp_path)
            if not reader.is_encrypted:
                return jsonify({
                    'error': 'PDF is not password protected',
                    'conversion_id': conversion_id
                }), 400

            if not reader.decrypt(password):
                return jsonify({
                    'error': 'Incorrect password',
                    'conversion_id': conversion_id
                }), 401

            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)

            # Save unlocked PDF
            with open(unlocked_path, "wb") as f:
                writer.write(f)

            return jsonify({
                'success': True,
                'filename': unlocked_filename,
                'download_url': url_for('download_unlocked', filename=unlocked_filename, _external=True),
                'conversion_id': conversion_id
            })

        except Exception as pdf_error:
            return jsonify({
                'error': f'PDF processing error: {str(pdf_error)}',
                'conversion_id': conversion_id
            }), 500

    except Exception as e:
        return jsonify({
            'error': f'Server error: {str(e)}',
            'conversion_id': conversion_id if 'conversion_id' in locals() else None
        }), 500

    finally:
        # Cleanup temporary files
        if 'temp_path' in locals() and os.path.exists(temp_path):
            os.remove(temp_path)

@app.route('/download-unlocked/<filename>')
def download_unlocked(filename):
    try:
        file_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)

        if not os.path.exists(file_path):
            return jsonify({'error': 'File not found'}), 404

        return send_from_directory(
            app.config['OUTPUT_FOLDER'],
            filename,
            as_attachment=True,
            mimetype='application/pdf',
            download_name=filename
        )
    except Exception as e:
        return jsonify({'error': f'Download error: {str(e)}'}), 500


# iWork to PDF
def allowed_iwork_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_IWORK_EXTENSIONS']

def convert_iwork_to_pdf(input_path, output_path):
    try:
        cmd = f"soffice --headless --convert-to pdf {input_path} --outdir {os.path.dirname(output_path)}"
        result = subprocess.run(cmd, shell=True, check=True, capture_output=True, text=True)
        print("Conversion output:", result.stdout)
        print("Conversion errors:", result.stderr)
        return True
    except subprocess.CalledProcessError as e:
        print(f"Conversion error: {e.stderr}")
        return False
    except Exception as e:
        print(f"Unexpected error: {str(e)}")
        return False

@app.route('/convert-iwork-to-pdf', methods=['POST'])
def handle_iwork_conversion():
    if 'file' not in request.files:
        return jsonify({'success': False, 'message': 'No file uploaded'}), 400

    file = request.files['file']
    user_id = request.form.get('user_id')
    conversion_id = str(uuid.uuid4())

    if file.filename == '':
        return jsonify({
            'success': False, 
            'message': 'No selected file',
            'conversion_id': conversion_id
        }), 400

    if not allowed_iwork_file(file.filename):
        return jsonify({
            'success': False, 
            'message': 'Invalid file type',
            'conversion_id': conversion_id
        }), 400

    try:
        filename = secure_filename(file.filename)
        unique_id = uuid.uuid4().hex
        upload_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{unique_id}_{filename}")
        file.save(upload_path)

        output_filename = f"{os.path.splitext(filename)[0]}.pdf"
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{unique_id}_{output_filename}")

        if convert_iwork_to_pdf(upload_path, output_path):
            pdf_url = f"/outputs/{unique_id}_{output_filename}"
            return jsonify({
                'success': True,
                'pdfUrl': pdf_url,
                'filename': output_filename,
                'conversion_id': conversion_id
            })
        
        return jsonify({
            'success': False, 
            'message': 'Conversion failed',
            'conversion_id': conversion_id
        }), 500
        
    except Exception as e:
        print(f"Server error: {str(e)}")
        return jsonify({
            'success': False, 
            'message': 'Server error during conversion',
            'conversion_id': conversion_id
        }), 500

@app.route('/outputs/<filename>')
def serve_iwork_pdf(filename):
    try:
        file_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
        if not os.path.exists(file_path):
            return jsonify({'success': False, 'message': 'File not found'}), 404
        
        # Send file with correct headers
        response = make_response(send_file(
            file_path,
            mimetype='application/pdf',
            as_attachment=False,
            download_name=filename
        ))
        
        # Critical headers to prevent caching issues
        response.headers['Cache-Control'] = 'no-store, no-cache, must-revalidate'
        response.headers['Pragma'] = 'no-cache'
        response.headers['Expires'] = '0'
        return response
        
    except Exception as e:
        print(f"Error serving PDF: {str(e)}")
        return jsonify({'success': False, 'message': 'Error serving PDF'}), 500


# eBook to PDF 
def allowed_ebook_file(filename):
    ALLOWED_EBOOK_EXTENSIONS = {'epub', 'mobi', 'azw', 'fb2'}
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EBOOK_EXTENSIONS

@app.route('/convert-ebook-to-pdf', methods=['POST'])
def handle_ebook_conversion():
    if 'file' not in request.files:
        return jsonify({'success': False, 'message': 'No file uploaded'}), 400

    file = request.files['file']
    user_id = request.form.get('user_id')
    conversion_id = str(uuid.uuid4())

    if file.filename == '':
        return jsonify({
            'success': False, 
            'message': 'No selected file',
            'conversion_id': conversion_id
        }), 400

    if not allowed_ebook_file(file.filename):
        return jsonify({
            'success': False, 
            'message': 'Invalid file type. Supported formats: EPUB, MOBI, AZW, FB2',
            'conversion_id': conversion_id
        }), 400

    try:
        filename = secure_filename(file.filename)
        unique_id = uuid.uuid4().hex
        upload_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{unique_id}_{filename}")
        output_filename = f"{os.path.splitext(filename)[0]}.pdf"
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{unique_id}_{output_filename}")

        os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
        os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

        file.save(upload_path)

        if os.path.getsize(upload_path) == 0:
            os.remove(upload_path)
            return jsonify({
                'success': False, 
                'message': 'Uploaded file is empty',
                'conversion_id': conversion_id
            }), 400

        cmd = [
            'ebook-convert',
            upload_path,
            output_path,
            '--embed-all-fonts',
            '--enable-heuristics',
            '--pdf-page-margin-left', '20',
            '--pdf-page-margin-right', '20',
            '--pdf-page-margin-top', '20',
            '--pdf-page-margin-bottom', '20',
            '--output-profile', 'tablet'  
        ]

        try:
            result = subprocess.run(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                errors='replace'
            )
            
            if result.returncode != 0:
                error_msg = result.stderr if result.stderr else "Conversion failed"
                return jsonify({
                    'success': False,
                    'message': f'Conversion failed: {error_msg}',
                    'conversion_id': conversion_id
                }), 500

            if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                return jsonify({
                    'success': False,
                    'message': 'Conversion completed but output file is empty',
                    'conversion_id': conversion_id
                }), 500

            pdf_url = f"/ebook-outputs/{unique_id}_{output_filename}"
            return jsonify({
                'success': True,
                'pdfUrl': pdf_url,
                'filename': output_filename,
                'conversion_id': conversion_id
            })

        except Exception as e:
            return jsonify({
                'success': False,
                'message': f'Conversion failed: {str(e)}',
                'conversion_id': conversion_id
            }), 500

    except Exception as e:
        return jsonify({
            'success': False,
            'message': f'Server error: {str(e)}',
            'conversion_id': conversion_id
        }), 500
    finally:
        if 'upload_path' in locals() and os.path.exists(upload_path):
            try:
                os.remove(upload_path)
            except Exception as e:
                print(f"Error cleaning up file: {str(e)}")

                
@app.route('/ebook-outputs/<path:filename>')
def serve_ebook_pdf(filename):
    """Serve converted eBook PDF files from the output directory"""
    try:
        return send_from_directory(
            app.config['OUTPUT_FOLDER'],
            filename, 
            as_attachment=False,
            mimetype='application/pdf'
        )
    except FileNotFoundError:
        return jsonify({'success': False, 'message': 'PDF file not found'}), 404


# Sign Document 
import comtypes.client
@app.route('/upload-signing-document', methods=['POST'])
def upload_signing_document():
    file = request.files.get('document_file')
    user_id = request.form.get('user_id')
    
    if not file:
        return jsonify({"error": "No file provided"}), 400

    filename = secure_filename(file.filename)
    file_ext = os.path.splitext(filename)[1].lower()
    file_size = os.fstat(file.stream.fileno()).st_size
    
    # Create a unique filename to avoid conflicts
    unique_id = str(uuid.uuid4())
    temp_filename = f"{unique_id}{file_ext}"
    temp_path = os.path.join(app.config['UPLOAD_FOLDER'], temp_filename)
    file.save(temp_path)

    # Record the initial document upload in database
    conversion_id = str(uuid.uuid4())
    try:
        db = get_db()
        cursor = db.cursor()
        cursor.execute("""
            INSERT INTO conversions 
            (user_id, conversion_type, original_filename, converted_filename, 
             file_size, status, conversion_id, parent_conversion_id) 
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        """, (
            user_id,
            'document_signing_upload',
            filename,
            None,  # Will be updated after signing
            file_size,
            'uploaded',  # Initial status
            conversion_id,
            None  # No parent for initial upload
        ))
        db.commit()
    except Exception as e:
        app.logger.error(f"Failed to record document upload: {str(e)}")
        return jsonify({"error": "Database error"}), 500

    # Convert Word to PDF if needed
    if file_ext in ['.doc', '.docx']:
        pdf_filename = f"{unique_id}.pdf"
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_filename)
        
        # Try multiple conversion methods with fallbacks
        conversion_success = False
        
        # Method 1: Try LibreOffice first (fastest)
        try:
            result = subprocess.run(
                ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', app.config['UPLOAD_FOLDER'], temp_path],
                capture_output=True,
                text=True,
                timeout=180
            )
            if result.returncode == 0:
                conversion_success = True
                # LibreOffice creates output with same name but .pdf extension
                temp_pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{unique_id}.pdf")
                if os.path.exists(temp_pdf_path):
                    os.rename(temp_pdf_path, pdf_path)
        except Exception as e:
            app.logger.error(f"LibreOffice conversion failed: {str(e)}")

        # Method 2: Try docx2pdf if LibreOffice failed
        if not conversion_success:
            try:
                from docx2pdf import convert
                convert(temp_path, pdf_path)
                conversion_success = True
            except Exception as e:
                app.logger.error(f"docx2pdf conversion failed: {str(e)}")

        # Method 3: Fallback to Word COM if other methods fail
        if not conversion_success:
            try:
                pythoncom.CoInitialize()
                word = comtypes.client.CreateObject('Word.Application')
                word.Visible = False
                doc = word.Documents.Open(temp_path)
                doc.SaveAs(pdf_path, FileFormat=17)  # 17 = wdFormatPDF
                doc.Close()
                word.Quit()
                conversion_success = True
            except Exception as e:
                app.logger.error(f"Word COM conversion failed: {str(e)}")
                # Update database with failure
                if user_id:
                    try:
                        cursor.execute("""
                            UPDATE conversions SET status=?, error_message=?
                            WHERE conversion_id=?
                        """, ('failed', str(e), conversion_id))
                        db.commit()
                    except Exception as db_error:
                        app.logger.error(f"Failed to update conversion record: {str(db_error)}")
                
                return jsonify({"error": "Failed to convert Word document"}), 500
            finally:
                pythoncom.CoUninitialize()

        # Clean up original Word file
        os.remove(temp_path)
        
        if not conversion_success:
            return jsonify({"error": "All conversion methods failed"}), 500
            
        return jsonify({
            "document_id": pdf_filename,
            "conversion_id": conversion_id
        })

    # Handle actual PDF uploads
    elif file_ext == '.pdf':
        # Rename to our unique filename
        final_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{unique_id}.pdf")
        os.rename(temp_path, final_path)
        return jsonify({
            "document_id": f"{unique_id}.pdf",
            "conversion_id": conversion_id
        })
    
    # Unsupported file type
    os.remove(temp_path)
    return jsonify({"error": "Unsupported file type"}), 400
    
@app.route('/sign-document/<document_id>')
def sign_document(document_id):
    document_path = os.path.join(app.config['UPLOAD_FOLDER'], document_id)
    if not os.path.exists(document_path):
        return "Document not found", 404

    # Get parameters from query string
    user_id = request.args.get('user_id')
    conversion_id = request.args.get('conversion_id')
    
    # Get original extension
    original_extension = os.path.splitext(document_id)[1].lower()
    
    # Determine which file to display (original or converted PDF)
    if original_extension in ('.doc', '.docx'):
        display_document_id = f"{os.path.splitext(document_id)[0]}.pdf"
    else:
        display_document_id = document_id

    
    return render_template_string("""
    <!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no">
  <title>Sign Document</title>
  <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.0/dist/css/bootstrap.min.css" rel="stylesheet">
  <script src="https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.4.120/pdf.min.js"></script>
  <script>
    pdfjsLib.GlobalWorkerOptions.workerSrc = 'https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.4.120/pdf.worker.min.js';
  </script>
  <style>
    body {
      padding-bottom: 20px;
      touch-action: manipulation;
    }
    #pdf-container {
      width: 100%;
      height: 60vh;
      overflow: auto;
      border: 1px solid #ddd;
      margin-bottom: 20px;
      background-color: #f5f5f5;
      position: relative;
      -webkit-overflow-scrolling: touch;
    }
    #pdfViewer {
      width: 100%;
      padding: 10px 0;
    }
    .page-container {
      margin: 0 auto 20px;
      box-shadow: 0 0 5px rgba(0,0,0,0.3);
      position: relative;
      background-color: white;
      max-width: 800px;
    }
    canvas.page {
      display: block;
      margin: 0 auto;
      width: 100% !important;
      height: auto !important;
    }
    .signature-preview {
      position: absolute;
      z-index: 100;
      max-width: 150px;
      border: 1px dashed #888;
      overflow: hidden;
      touch-action: none;
    }
    #signatureCanvas {
      border: 1px solid #000;
      width: 100%;
      height: 150px;
      background-color: white;
      touch-action: none;
    }
    #typedPreview {
      font-family: 'Brush Script MT', cursive;
      font-size: 1.5rem;
      min-height: 40px;
      border: 1px dashed #ccc;
      padding: 8px;
      margin-top: 8px;
    }
    #uploadPreview {
      max-height: 120px;
      display: none;
      margin-top: 8px;
    }
    .instructions {
      background-color: #f8f9fa;
      padding: 12px;
      border-radius: 5px;
      margin-bottom: 15px;
      font-size: 0.9rem;
    }
    .resize-handle {
      position: absolute;
      width: 24px;
      height: 24px;
      background: #555;
      right: 0;
      bottom: 0;
      cursor: nwse-resize;
      z-index: 101;
    }
    .close-signature {
      position: absolute;
      top: 0;
      right: 0;
      background: rgba(255,255,255,0.7);
      border: 1px solid #aaa;
      cursor: pointer;
      padding: 2px 8px;
      font-size: 18px;
      font-weight: bold;
      z-index: 102;
      width: 28px;
      height: 28px;
      display: flex;
      align-items: center;
      justify-content: center;
    }
    .btn {
      padding: 10px 16px;
      font-size: 1rem;
    }
    .nav-tabs .nav-link {
      padding: 10px 12px;
    }
    .tab-content {
      padding: 15px 0;
    }
    .loading-message {
      display: flex;
      justify-content: center;
      align-items: center;
      height: 100%;
      flex-direction: column;
    }
    .spinner {
      width: 3rem;
      height: 3rem;
      margin-bottom: 1rem;
    }
    /* Improved touch targets */
    .btn, .nav-link, .close-signature, .resize-handle {
      touch-action: manipulation;
    }
  </style>
</head>
<body>
<div class="container mt-3">
  <h1 class="h4">Sign Document</h1>
  <div class="instructions alert alert-info py-2">
    <strong>Instructions:</strong> Tap anywhere on the document to place your signature. 
    Drag to reposition and pinch to resize. When finished, tap "Save & Download".
  </div>
  
  <!-- PDF Viewer Container -->
  <div id="pdf-container">
    <div id="pdfViewer">
      <div class="loading-message">
        <div class="spinner-border text-primary spinner" role="status">
          <span class="visually-hidden">Loading...</span>
        </div>
        <p>Loading document...</p>
      </div>
    </div>
  </div>
  
  <button id="saveDocument" class="btn btn-success mt-2 w-100">Save & Download Signed Document</button>
</div>

<!-- Signature Modal -->
<div class="modal fade" id="signatureModal" tabindex="-1" aria-hidden="true">
  <div class="modal-dialog modal-fullscreen-sm-down">
    <div class="modal-content">
      <div class="modal-header">
        <h5 class="modal-title">Create Signature</h5>
        <button type="button" class="btn-close" data-bs-dismiss="modal" aria-label="Close"></button>
      </div>
      <div class="modal-body">
        <ul class="nav nav-tabs">
          <li class="nav-item"><button class="nav-link active" data-bs-toggle="tab" data-bs-target="#draw">Draw</button></li>
          <li class="nav-item"><button class="nav-link" data-bs-toggle="tab" data-bs-target="#type">Type</button></li>
          <li class="nav-item"><button class="nav-link" data-bs-toggle="tab" data-bs-target="#upload">Upload</button></li>
        </ul>

        <div class="tab-content mt-3">
          <div class="tab-pane fade show active" id="draw">
            <canvas id="signatureCanvas"></canvas>
            <button id="clearDraw" class="btn btn-warning mt-2 w-100">Clear</button>
          </div>

          <div class="tab-pane fade" id="type">
            <input type="text" id="typedSignature" class="form-control form-control-lg" placeholder="Type your name">
            <div id="typedPreview">Signature will appear here</div>
            <button id="clearType" class="btn btn-warning mt-2 w-100">Clear</button>
          </div>

          <div class="tab-pane fade" id="upload">
            <input type="file" id="uploadSignature" class="form-control form-control-lg" accept="image/*">
            <img id="uploadPreview" class="img-fluid mt-2">
            <button id="clearUpload" class="btn btn-warning mt-2 w-100">Clear</button>
          </div>
        </div>
      </div>
      <div class="modal-footer">
        <button id="addSignature" class="btn btn-primary w-100">Add Signature</button>
        <button type="button" class="btn btn-secondary w-100 mt-2" data-bs-dismiss="modal">Close</button>
      </div>
    </div>
  </div>
</div>

<script src="https://cdn.jsdelivr.net/npm/bootstrap@5.3.0/dist/js/bootstrap.bundle.min.js"></script>
<script>
// Configuration
const backendUrl = "http://localhost:5000"; // ADDED BACKEND URL
const pdfUrl = `${backendUrl}/uploads/{{ document_id }}`; // MODIFIED TO USE BACKEND URL
const pdfViewer = document.getElementById("pdfViewer");
const originalExtension = "{{ original_extension }}";
let currentPage = 0;
let pageContainers = [];
let signatures = [];
let pdfDoc = null;
let pageRects = [];
let clickX = 0;
let clickY = 0;
let signatureCanvas, signatureCtx; // Global canvas reference

// [Rest of the initialization code remains exactly the same...]

// Save signed document
async function saveSignedDocument() {
  if (signatures.length === 0) {
    alert("Please add at least one signature first");
    return;
  }

  // Show loading state
  const saveButton = document.getElementById("saveDocument");
  const originalText = saveButton.textContent;
  saveButton.disabled = true;
  saveButton.innerHTML = '<span class="spinner-border spinner-border-sm" role="status" aria-hidden="true"></span> Processing...';

  try {
    // Prepare signature data
    const signatureData = signatures.map(sig => {
      const pageContainer = sig.element.parentElement;
      const canvas = pageContainer.querySelector('canvas');

      // ****** ONLY CHANGE: convert from displayed CSS pixels to original PDF pixels ******
      const pageInfo = pageRects[sig.page - 1];
      const canvasRect = canvas.getBoundingClientRect();

      // ratios from what's displayed on screen to the PDF's original size
      const ratioX = pageInfo.viewport.width / canvasRect.width;
      const ratioY = pageInfo.viewport.height / canvasRect.height;

      const originalLeft = sig.position.left * ratioX;
      const originalTop = sig.position.top * ratioY;
      const originalWidth = sig.position.width * ratioX;
      const originalHeight = sig.position.height * ratioY;
      // ****************************************************************************************

      return {
        page: sig.page,
        signature: sig.data,
        position: {
          x: originalLeft,
          y: pageInfo.viewport.height - originalTop - originalHeight, // invert Y for PDF space
          width: originalWidth,
          height: originalHeight,
          originalWidth: pageInfo.viewport.width,
          originalHeight: pageInfo.viewport.height
        }
      };
    });

    const response = await fetch(`${backendUrl}/save-signed-document`, { // MODIFIED TO USE BACKEND URL
      method: "POST",
      headers: { 
        "Content-Type": "application/json",
      },
      body: JSON.stringify({
        document_id: "{{ document_id }}",
        original_extension: originalExtension,
        signatures: signatureData,
        user_id: "{{ user_id }}",
        conversion_id: "{{ conversion_id }}"
      })
    });

    if (!response.ok) {
      throw new Error(`Server returned ${response.status}: ${response.statusText}`);
    }

    const blob = await response.blob();
    const a = document.createElement('a');
    a.href = window.URL.createObjectURL(blob);
    a.download = `signed_{{ document_id }}`;
    document.body.appendChild(a);
    a.click();
    a.remove();
  } catch (error) {
    console.error("Error saving document:", error);
    alert("Error saving document: " + error.message);
  } finally {
    saveButton.disabled = false;
    saveButton.textContent = originalText;
  }
}
</script>
</body>
</html>
    """,    document_id=display_document_id, 
    original_extension=original_extension,
    user_id=user_id,
    conversion_id=conversion_id)


@app.route('/save-signed-document', methods=['POST'])
def save_signed_document():
    data = request.json
    document_id = data.get('document_id')
    original_extension = data.get('original_extension')
    signatures_data = data.get('signatures', [])
    user_id = data.get('user_id')
    conversion_id = data.get('conversion_id')  # This is the parent conversion_id

    if not document_id or not signatures_data:
        return jsonify({"error": "Invalid data"}), 400

    document_path = os.path.join(app.config['UPLOAD_FOLDER'], document_id)
    if not os.path.exists(document_path):
        return jsonify({"error": "Document not found"}), 404

    try:
        # Process all signatures
        signature_images = []
        for sig in signatures_data:
            if sig['signature'].startswith('data:'):
                signature_image = base64.b64decode(sig['signature'].split(',')[1])
            else:
                import requests
                response = requests.get(sig['signature'])
                signature_image = response.content

            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_sig:
                tmp_sig.write(signature_image)
                signature_images.append({
                    'path': tmp_sig.name,
                    'page': sig['page'] - 1,  # Convert to 0-based index
                    'position': sig['position']
                })

        # Generate a new conversion_id for the signed version
        signed_conversion_id = str(uuid.uuid4())
        signed_filename = f"signed_{document_id}"
        signed_document_path = os.path.join(app.config['OUTPUT_FOLDER'], signed_filename)

        # Handle PDF documents
        original_pdf = PdfReader(document_path)
        output = PdfWriter()

        for page_num in range(len(original_pdf.pages)):
            page = original_pdf.pages[page_num]
            page_width = float(page.mediabox.width)
            page_height = float(page.mediabox.height)

            page_signatures = [sig for sig in signature_images if sig['page'] == page_num]
            
            if page_signatures:
                packet = BytesIO()
                can = canvas.Canvas(packet, pagesize=(page_width, page_height))
                
                for sig in page_signatures:
                    pos = sig['position']
                    x = pos['x']
                    y = pos['y']
                    width = pos['width']
                    height = pos['height']
                    can.drawImage(sig['path'], x, y, width, height, mask='auto')
                
                can.save()
                packet.seek(0)
                overlay_pdf = PdfReader(packet)
                page.merge_page(overlay_pdf.pages[0])
            
            output.add_page(page)

        with open(signed_document_path, "wb") as output_pdf:
            output.write(output_pdf)

        # Update database with completion - create a new record for the signed version
        if user_id and conversion_id:
            try:
                db = get_db()
                cursor = db.cursor()
                
                # First get the original conversion details
                cursor.execute("""
                    SELECT original_filename, file_size FROM conversions 
                    WHERE conversion_id=?
                """, (conversion_id,))
                original_data = cursor.fetchone()
                
                if original_data:
                    # Create new record for the signed document
                    cursor.execute("""
                        INSERT INTO conversions 
                        (user_id, conversion_type, original_filename, converted_filename, 
                         file_size, status, conversion_id, parent_conversion_id, completed_at) 
                        VALUES (?, ?, ?, ?, ?, ?, ?, ?, CURRENT_TIMESTAMP)
                    """, (
                        user_id,
                        'document_signing_completed',
                        original_data['original_filename'],
                        signed_filename,
                        os.path.getsize(signed_document_path),
                        'completed',
                        signed_conversion_id,
                        conversion_id  # Link to parent conversion
                    ))
                    
                    # Update the original record to mark it as processed
                    cursor.execute("""
                        UPDATE conversions SET 
                        status=?,
                        completed_at=CURRENT_TIMESTAMP
                        WHERE conversion_id=?
                    """, (
                        'processed',
                        conversion_id
                    ))
                    
                    db.commit()
            except Exception as e:
                app.logger.error(f"Failed to update conversion records: {str(e)}")

        # Clean up temporary files
        for sig in signature_images:
            os.unlink(sig['path'])

        return send_from_directory(app.config['OUTPUT_FOLDER'], os.path.basename(signed_document_path), as_attachment=True)

    except Exception as e:
        # Update database with failure if we have a conversion_id
        if user_id and conversion_id:
            try:
                db = get_db()
                cursor = db.cursor()
                cursor.execute("""
                    UPDATE conversions SET 
                    status=?, 
                    error_message=?,
                    completed_at=CURRENT_TIMESTAMP
                    WHERE conversion_id=?
                """, (
                    'failed',
                    str(e),
                    conversion_id
                ))
                db.commit()
            except Exception as db_error:
                app.logger.error(f"Failed to update conversion record: {str(db_error)}")
        
        return jsonify({"error": str(e)}), 500

@app.route('/download-signed/<filename>')
def download_signed(filename):
    return send_from_directory(app.config['OUTPUT_FOLDER'], filename, as_attachment=True)



# Edit PDF document 
@app.route('/upload-edit-pdf', methods=['POST'])
def upload_edit_pdf():
    if 'pdfFile' not in request.files:
        return jsonify({'error': 'No file part'}), 400

    file = request.files['pdfFile']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    if not file.filename.lower().endswith('.pdf'):
        return jsonify({'error': 'Only PDF files are allowed'}), 400

    try:
        # Get current user ID if available
        current_user_id = request.form.get('user_id') or 'anonymous'
        
        filename = secure_filename(file.filename)
        unique_id = str(uuid.uuid4())
        pdf_filename = f"{unique_id}_{filename}"
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_filename)
        file.save(pdf_path)

        docx_filename = f"{unique_id}.docx"
        docx_path = os.path.join(app.config['UPLOAD_FOLDER'], docx_filename)
        
        # Record initial upload in database
        conversion_id = str(uuid.uuid4())
        db = get_db()
        cursor = db.cursor()
        cursor.execute("""
            INSERT INTO conversions 
            (user_id, conversion_type, original_filename, converted_filename, file_size, status, conversion_id) 
            VALUES (?, ?, ?, ?, ?, ?, ?)
        """, (
            current_user_id,
            'pdf_edit_upload',
            file.filename,
            pdf_filename,
            os.path.getsize(pdf_path),
            'uploaded',
            conversion_id
        ))
        db.commit()
        
        # Convert PDF to DOCX with enhanced error handling
        try:
            converter = Converter(pdf_path)
            try:
                converter.convert(docx_path, keep_images=True)
            except Exception as img_error:
                app.logger.warning(f"Image conversion failed, trying without images: {str(img_error)}")
                converter.convert(docx_path, keep_images=False)
            finally:
                converter.close()
        except Exception as conv_error:
            # Update database with failure
            cursor.execute("""
                UPDATE conversions SET status = ?, error_message = ?
                WHERE conversion_id = ?
            """, ('failed', str(conv_error), conversion_id))
            db.commit()
            
            if os.path.exists(pdf_path):
                os.remove(pdf_path)
            return jsonify({'error': f'PDF conversion failed: {str(conv_error)}'}), 500

        # Process DOCX to HTML with images and styles
        try:
            doc = Document(docx_path)
            html_content = process_docx_to_html(doc, unique_id)
            
            # Update database with success
            cursor.execute("""
                UPDATE conversions SET status = ?, converted_filename = ?
                WHERE conversion_id = ?
            """, ('ready_for_edit', docx_filename, conversion_id))
            db.commit()
            
        except Exception as proc_error:
            cursor.execute("""
                UPDATE conversions SET status = ?, error_message = ?
                WHERE conversion_id = ?
            """, ('failed', str(proc_error), conversion_id))
            db.commit()
            return jsonify({'error': f'DOCX processing failed: {str(proc_error)}'}), 500

        return jsonify({
            'success': True,
            'pdfFileName': pdf_filename,
            'docxFilename': docx_filename,
            'htmlContent': html_content,
            'conversionId': conversion_id
        })

    except Exception as e:
        app.logger.error(f"Error in upload-edit-pdf: {str(e)}\n{traceback.format_exc()}")
        if 'pdf_path' in locals() and os.path.exists(pdf_path):
            os.remove(pdf_path)
        if 'docx_path' in locals() and os.path.exists(docx_path):
            os.remove(docx_path)
        return jsonify({'error': 'Internal server error'}), 500

@app.route('/get-pdf-content-for-editing')
def get_pdf_content_for_editing():
    try:
        file_id = request.args.get('file_id')
        if not file_id:
            return jsonify({'error': 'Missing file_id parameter'}), 400

        # Clean the filename and extract UUID
        clean_file_id = secure_filename(file_id)
        unique_id = clean_file_id.split('_')[0]
        docx_filename = f"{unique_id}.docx"
        docx_path = os.path.join(app.config['UPLOAD_FOLDER'], docx_filename)
        
        # Verify the file exists
        if not os.path.exists(docx_path):
            existing_files = os.listdir(app.config['UPLOAD_FOLDER'])
            return jsonify({
                'error': f'Document not found at path: {docx_path}',
                'looking_for': docx_filename,
                'existing_files': existing_files
            }), 404

        # Process the DOCX file
        doc = Document(docx_path)
        html_content = process_docx_to_html(doc, unique_id)

        return jsonify({'success': True, 'content': html_content})
    except Exception as e:
        app.logger.error(f"Error in get-pdf-content: {str(e)}\n{traceback.format_exc()}")
        return jsonify({'error': 'Internal server error'}), 500

@app.route('/pdf-editor-view')
def pdf_editor_view():
    file_id = request.args.get('file')
    if not file_id:
        return "Missing file parameter", 400
    return render_template('editor.html', file_id=file_id)

def process_docx_to_html(doc, unique_id):
    """Convert DOCX to HTML while preserving images and styles"""
    html_parts = []
    images_dir = os.path.join(app.config['UPLOAD_FOLDER'], f"{unique_id}_images")
    os.makedirs(images_dir, exist_ok=True)
    
    for para in doc.paragraphs:
        if para.text.strip():
            # Process runs to preserve styles
            para_html = []
            for run in para.runs:
                text = run.text
                if not text.strip():
                    continue
                    
                style = []
                if run.bold:
                    style.append('font-weight:bold')
                if run.italic:
                    style.append('font-style:italic')
                if run.underline:
                    style.append('text-decoration:underline')
                if run.font.color.rgb:
                    style.append(f'color:#{run.font.color.rgb[2:]}')
                    
                style_str = ';'.join(style)
                para_html.append(f'<span style="{style_str}">{text}</span>')
            
            if para_html:
                html_parts.append(f'<p>{"".join(para_html)}</p>')
    
    # Process images
    rels = doc.part.rels
    for rel in rels:
        if "image" in rels[rel].target_ref:
            image_part = rels[rel].target_part
            image_ext = image_part.content_type.split('/')[-1]
            image_filename = f"{uuid.uuid4()}.{image_ext}"
            image_path = os.path.join(images_dir, image_filename)
            
            with open(image_path, 'wb') as f:
                f.write(image_part.blob)
            
            # Convert image to base64 for HTML embedding
            with open(image_path, 'rb') as img_file:
                img_data = base64.b64encode(img_file.read()).decode('utf-8')
                img_src = f"data:image/{image_ext};base64,{img_data}"
                html_parts.append(f'<img src="{img_src}" style="max-width:100%;" data-resized="true">')

    return "".join(html_parts)


@app.route('/save-edited-pdf', methods=['POST'])
def save_edited_pdf():
    try:
        data = request.get_json()
        if not data:
            return jsonify({'error': 'No JSON data received'}), 400

        file_id = data.get('file_id')
        html_content = data.get('content')
        conversion_id = data.get('conversionId')
        current_user_id = data.get('userId') or 'anonymous'
        
        if not file_id or not html_content:
            return jsonify({'error': 'Missing parameters'}), 400

        # Clean the filename and extract UUID
        clean_file_id = secure_filename(file_id)
        unique_id = clean_file_id.split('_')[0]
        
        # Create BeautifulSoup object and remove duplicates
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Track seen images and remove duplicates
        seen_images = {}
        for img in soup.find_all('img'):
            src = img.get('src', '')
            if src in seen_images:
                img.decompose()
            else:
                img['data-resized'] = 'true'
                seen_images[src] = img
                if 'style' in img.attrs:
                    style = img['style']
                    width_match = re.search(r'width:\s*(\d+)px', style)
                    if width_match:
                        img['data-width'] = width_match.group(1)
        
        cleaned_html = str(soup)
        
        # Create temporary directory for processing
        temp_dir = os.path.join(app.config['UPLOAD_FOLDER'], f"temp_{unique_id}")
        os.makedirs(temp_dir, exist_ok=True)
        
        temp_docx_path = os.path.join(temp_dir, f"{unique_id}.docx")
        temp_pdf_path = os.path.join(temp_dir, clean_file_id)
        
        # Create images directory
        images_dir = os.path.join(temp_dir, f"{unique_id}_images")
        os.makedirs(images_dir, exist_ok=True)

        # Create a new document with styles
        doc = Document()
        
        # Add styles to the document
        styles = doc.styles
        if 'Heading 1' not in styles:
            styles.add_style('Heading 1', WD_STYLE_TYPE.PARAGRAPH)
        if 'Heading 2' not in styles:
            styles.add_style('Heading 2', WD_STYLE_TYPE.PARAGRAPH)
        
        # Parse cleaned HTML with BeautifulSoup
        soup = BeautifulSoup(cleaned_html, 'html.parser')
        
        # Process all elements in order
        for element in soup.find_all(True):
            if element.name == 'p':
                para = doc.add_paragraph()
                process_element(element, para, images_dir)
            elif element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                heading_level = int(element.name[1])
                para = doc.add_paragraph(style=f'Heading {heading_level}')
                process_element(element, para, images_dir)
            elif element.name == 'img':
                if element.get('data-resized') == 'true':
                    try:
                        width = None
                        if 'data-width' in element.attrs:
                            width = int(element['data-width'])
                        elif 'style' in element.attrs:
                            style = element['style']
                            width_match = re.search(r'width:\s*(\d+)px', style)
                            if width_match:
                                width = int(width_match.group(1))
                        add_image_to_doc(element, doc, images_dir, width)
                    except Exception as img_error:
                        app.logger.error(f"Failed to process image: {str(img_error)}")
                        continue  # Skip this image but continue processing
            elif element.name in ['ul', 'ol']:
                process_list(element, doc, images_dir, element.name == 'ol')
            elif element.name == 'table':
                process_table(element, doc, images_dir)

        # Save the DOCX file
        try:
            doc.save(temp_docx_path)
        except Exception as e:
            temp_docx_path = os.path.join(temp_dir, f"{unique_id}_alt.docx")
            doc.save(temp_docx_path)

        # Convert DOCX to PDF
        final_pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], clean_file_id)
        success = False
        error_msg = None
        
        # Method 1: Try LibreOffice
        try:
            import subprocess
            result = subprocess.run(
                ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', temp_dir, temp_docx_path],
                capture_output=True,
                text=True,
                timeout=180
            )
            if result.returncode == 0:
                success = True
        except Exception as e:
            error_msg = f"LibreOffice failed: {str(e)}"

        # Method 2: Fallback to Word COM
        if not success:
            try:
                pythoncom.CoInitialize()
                word = Dispatch('Word.Application')
                word.Visible = False
                word.DisplayAlerts = False
                
                doc = word.Documents.Open(os.path.abspath(temp_docx_path))
                doc.SaveAs(os.path.abspath(temp_pdf_path), FileFormat=17)
                doc.Close(False)
                word.Quit()
                success = True
            except Exception as e:
                error_msg = f"Word COM failed: {str(e)}"
            finally:
                pythoncom.CoUninitialize()

        # Method 3: Final fallback to docx2pdf
        if not success:
            try:
                from docx2pdf import convert
                convert(temp_docx_path, temp_pdf_path)
                success = True
            except Exception as e:
                error_msg = f"docx2pdf failed: {str(e)}"

        if not success:
            raise Exception(f"All PDF conversion methods failed. Last error: {error_msg}")

        # Verify PDF was created
        if not os.path.exists(temp_pdf_path):
            temp_pdf_path = os.path.join(temp_dir, f"{unique_id}.pdf")
            if not os.path.exists(temp_pdf_path):
                raise Exception("PDF file was not created")

        # Move the final PDF
        shutil.move(temp_pdf_path, final_pdf_path)
        
        # After successful PDF creation, record in database
        db = get_db()
        cursor = db.cursor()
        
        # Create a new record for the edited version
        edited_conversion_id = str(uuid.uuid4())
        cursor.execute("""
            INSERT INTO conversions 
            (user_id, conversion_type, original_filename, converted_filename, 
             file_size, status, conversion_id, parent_conversion_id) 
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        """, (
            current_user_id,
            'pdf_edit_save',
            clean_file_id,
            os.path.basename(final_pdf_path),
            os.path.getsize(final_pdf_path),
            'completed',
            edited_conversion_id,
            conversion_id
        ))
        db.commit()

        # Clean up temporary files
        try:
            shutil.rmtree(temp_dir)
        except Exception as e:
            app.logger.error(f"Error cleaning up temp files: {str(e)}")

        return jsonify({
            'success': True, 
            'pdfUrl': os.path.basename(final_pdf_path),
            'message': 'PDF successfully saved',
            'conversionId': edited_conversion_id
        })
        
    except Exception as e:
        app.logger.error(f"Error saving PDF: {str(e)}\n{traceback.format_exc()}")
        # Record failure if we had a conversion_id
        if 'conversion_id' in locals():
            try:
                db = get_db()
                cursor = db.cursor()
                cursor.execute("""
                    UPDATE conversions SET status = ?, error_message = ?
                    WHERE conversion_id = ?
                """, ('failed', str(e), conversion_id))
                db.commit()
            except Exception as db_error:
                app.logger.error(f"Failed to record error: {str(db_error)}")
                
        if 'temp_dir' in locals() and os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
            except Exception as cleanup_error:
                app.logger.error(f"Cleanup error: {str(cleanup_error)}")
        return jsonify({
            'error': str(e),
            'message': 'Failed to save PDF'
        }), 500



def process_element(element, para, images_dir):
    """Process HTML element with comprehensive style and content handling"""
    for content in element.contents:
        if content.name == 'img':
            # Get dimensions from style or attributes
            width = None
            if 'style' in content.attrs:
                style = content['style']
                width_match = re.search(r'width:\s*(\d+)px', style)
                if width_match:
                    width = int(width_match.group(1))
            
            add_image_to_paragraph(content, para, images_dir, width)
        elif content.name == 'span':
            run = para.add_run(content.text)
            apply_styles_to_run(run, content.get('style', ''))
        elif content.name == 'b' or content.name == 'strong':
            run = para.add_run(content.text)
            run.bold = True
        elif content.name == 'i' or content.name == 'em':
            run = para.add_run(content.text)
            run.italic = True
        elif content.name == 'u':
            run = para.add_run(content.text)
            run.underline = True
        elif content.name == 'br':
            para.add_run().add_break()
        elif content.name is None:  # Text node
            text = str(content).replace('\n', ' ')
            if text.strip():
                para.add_run(text)

def apply_styles_to_run(run, style_str):
    """Apply CSS styles to a docx run with comprehensive style support"""
    if not style_str:
        return
        
    styles = [s.strip() for s in style_str.split(';') if s.strip()]
    style_dict = {}
    
    for style in styles:
        if ':' not in style:
            continue
        prop, value = style.split(':', 1)
        prop = prop.strip().lower()
        value = value.strip().lower()
        style_dict[prop] = value
    
    # Apply styles
    if 'font-weight' in style_dict and style_dict['font-weight'] in ['bold', 'bolder', '700', '800', '900']:
        run.bold = True
    if 'font-style' in style_dict and style_dict['font-style'] == 'italic':
        run.italic = True
    if 'text-decoration' in style_dict and 'underline' in style_dict['text-decoration']:
        run.underline = True
    if 'color' in style_dict:
        color = style_dict['color']
        if color.startswith('#'):
            try:
                run.font.color.rgb = RGBColor.from_string(color[1:])
            except:
                pass
    if 'font-size' in style_dict:
        try:
            size = float(style_dict['font-size'].replace('pt', ''))
            run.font.size = Pt(size)
        except:
            pass

def process_image_element(element, parent, images_dir):
    """Process image element with better format support and error handling"""
    try:
        if 'src' not in element.attrs:
            return
            
        src = element['src']
        if not src.startswith('data:image'):
            return
            
        # Extract image format from data URL
        img_format = src.split(';')[0].split('/')[-1]
        if img_format not in ['png', 'jpeg', 'jpg', 'gif']:
            img_format = 'png'  # default to png if format not recognized
            
        img_data = src.split('base64,')[-1]
        img_bytes = base64.b64decode(img_data)
        img_filename = f"{uuid.uuid4()}.{img_format}"
        img_path = os.path.join(images_dir, img_filename)
        
        # Ensure directory exists
        os.makedirs(images_dir, exist_ok=True)
        
        with open(img_path, 'wb') as f:
            f.write(img_bytes)
        
        # Add image to document with size control
        if isinstance(parent, Document):
            parent.add_picture(img_path, width=Inches(6))  # Limit width to 6 inches
        else:  # Paragraph
            parent.add_run().add_picture(img_path, width=Inches(6))
            
    except Exception as e:
        app.logger.error(f"Error processing image: {str(e)}\n{traceback.format_exc()}")
        raise

def add_image_to_doc(element, doc, images_dir, width=None, height=None):
    """Add image directly to document with optional width/height in pixels"""
    try:
        if 'src' not in element.attrs:
            return
            
        src = element['src']
        if not src.startswith('data:image'):
            return
            
        # Skip if this image has already been processed
        if hasattr(element, '_processed'):
            return
        element._processed = True
            
        # Extract image format from data URL
        img_format = src.split(';')[0].split('/')[-1].lower()
        if img_format not in ['png', 'jpeg', 'jpg', 'gif']:
            img_format = 'png'  # default to png if format not recognized
            
        img_data = src.split('base64,')[-1]
        try:
            img_bytes = base64.b64decode(img_data)
            img_filename = f"{uuid.uuid4()}.{img_format}"
            img_path = os.path.join(images_dir, img_filename)
            
            with open(img_path, 'wb') as f:
                f.write(img_bytes)
            
            # Convert pixel dimensions to inches (96 DPI standard)
            width_inches = Inches(width / 96) if width else None
            height_inches = Inches(height / 96) if height else None
            
            # Add image with dimensions if specified
            try:
                if width_inches and height_inches:
                    doc.add_picture(img_path, width=width_inches, height=height_inches)
                elif width_inches:
                    doc.add_picture(img_path, width=width_inches)
                elif height_inches:
                    doc.add_picture(img_path, height=height_inches)
                else:
                    doc.add_picture(img_path)
            except Exception as e:
                app.logger.error(f"Failed to add picture to doc: {str(e)}")
                # Try adding without dimensions if that failed
                doc.add_picture(img_path)
                
        except Exception as e:
            app.logger.error(f"Error processing image data: {str(e)}")
            raise

    except Exception as e:
        app.logger.error(f"Error adding image to doc: {str(e)}")
        raise

def add_image_to_paragraph(element, para, images_dir, width=None):
    """Embed an image inside a paragraph with optional width in pixels"""
    try:
        if 'src' not in element.attrs or not element['src'].startswith('data:image'):
            return

        # Skip if this is not a resized image
        if element.get('data-resized') != 'true':
            return

        # Extract image format from data URL
        img_format = element['src'].split(';')[0].split('/')[-1]
        if img_format not in ['png', 'jpeg', 'jpg', 'gif']:
            img_format = 'png'
            
        img_data = element['src'].split('base64,')[-1]
        img_bytes = base64.b64decode(img_data)
        img_filename = f"{uuid.uuid4()}.{img_format}"
        img_path = os.path.join(images_dir, img_filename)

        with open(img_path, 'wb') as f:
            f.write(img_bytes)

        run = para.add_run()
        
        # Convert pixel width to inches if specified
        width_inches = Inches(width / 96) if width else None
        
        if width_inches:
            run.add_picture(img_path, width=width_inches)
        else:
            run.add_picture(img_path)
    except Exception as e:
        app.logger.error(f"Failed to embed image in paragraph: {str(e)}")
        raise


def add_image_to_paragraph(element, para, images_dir, width=None):
    """Embed an image inside a paragraph with optional width in pixels"""
    try:
        if 'src' not in element.attrs or not element['src'].startswith('data:image'):
            return

        # Extract image format from data URL
        img_format = element['src'].split(';')[0].split('/')[-1]
        if img_format not in ['png', 'jpeg', 'jpg', 'gif']:
            img_format = 'png'
            
        img_data = element['src'].split('base64,')[-1]
        img_bytes = base64.b64decode(img_data)
        img_filename = f"{uuid.uuid4()}.{img_format}"
        img_path = os.path.join(images_dir, img_filename)

        with open(img_path, 'wb') as f:
            f.write(img_bytes)

        run = para.add_run()
        
        # Convert pixel width to inches if specified
        width_inches = Inches(width / 96) if width else None
        
        if width_inches:
            run.add_picture(img_path, width=width_inches)
        else:
            run.add_picture(img_path)
    except Exception as e:
        app.logger.error(f"Failed to embed image in paragraph: {str(e)}")
        raise

def process_image_element(element, parent, images_dir):
    """Process image element"""
    try:
        if 'src' not in element.attrs:
            return
            
        src = element['src']
        if not src.startswith('data:image'):
            return
            
        img_data = src.split('base64,')[-1]
        img_bytes = base64.b64decode(img_data)
        img_filename = f"{uuid.uuid4()}.png"
        img_path = os.path.join(images_dir, img_filename)
        
        with open(img_path, 'wb') as f:
            f.write(img_bytes)
        
        if isinstance(parent, Document):
            parent.add_picture(img_path)
        else:  # Paragraph
            parent.add_run().add_picture(img_path)
    except Exception as e:
        app.logger.error(f"Error processing image: {str(e)}")

def process_list(element, doc, images_dir, is_ordered=False):
    """Process list elements (ul/ol)"""
    for item in element.find_all('li', recursive=False):
        if is_ordered:
            para = doc.add_paragraph(style='List Number')
        else:
            para = doc.add_paragraph(style='List Bullet')
        process_element_content(item, para, images_dir)

def process_table(element, doc, images_dir):
    """Process table elements with proper structure and images"""
    # Count rows and columns
    rows = element.find_all('tr')
    if not rows:
        return
        
    # Determine number of columns from first row
    cols = len(rows[0].find_all(['td', 'th']))
    
    table = doc.add_table(rows=1, cols=cols)
    
    for row_idx, row in enumerate(rows):
        cells = row.find_all(['td', 'th'])
        if not cells:
            continue
            
        # Add new row if needed (first row already exists)
        if row_idx > 0:
            table.add_row()
            
        for col_idx, cell in enumerate(cells):
            if col_idx >= len(table.rows[row_idx].cells):
                table.add_column()
                
            process_element_content(cell, table.rows[row_idx].cells[col_idx], images_dir)


def process_html_to_docx(html_content, doc, unique_id):
    """Improved HTML to DOCX conversion with better style handling including image resizing"""
    soup = BeautifulSoup(html_content, 'html.parser')
    images_dir = os.path.join(app.config['UPLOAD_FOLDER'], f"{unique_id}_images")
    os.makedirs(images_dir, exist_ok=True)
    
    # Process all elements in order
    for element in soup.find_all(True):
        if element.name == 'p':
            para = doc.add_paragraph()
            process_element_content(element, para, images_dir)
        elif element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
            heading_level = int(element.name[1])
            para = doc.add_paragraph(style=f'Heading {heading_level}')
            process_element_content(element, para, images_dir)
        elif element.name == 'img':
            # Handle standalone images with width from style or data attributes
            width = None
            height = None
            
            # Check style attribute first
            if 'style' in element.attrs:
                style = element['style']
                width_match = re.search(r'width:\s*(\d+)px', style)
                height_match = re.search(r'height:\s*(\d+)px', style)
                if width_match:
                    width = int(width_match.group(1))
                if height_match:
                    height = int(height_match.group(1))
            
            # Fallback to data attributes if no style
            if not width and 'data-width' in element.attrs:
                width = int(element['data-width'])
            if not height and 'data-height' in element.attrs:
                height = int(element['data-height'])
            
            add_image_to_doc(element, doc, images_dir, width, height)
        elif element.name in ['ul', 'ol']:
            process_list(element, doc, images_dir, element.name == 'ol')
        elif element.name == 'table':
            process_table(element, doc, images_dir)
        elif element.name == 'div':
            # Handle div containers (common in summernote output)
            for child in element.children:
                if child.name == 'img':
                    width = None
                    if 'style' in child.attrs:
                        style = child['style']
                        width_match = re.search(r'width:\s*(\d+)px', style)
                        if width_match:
                            width = int(width_match.group(1))
                    add_image_to_doc(child, doc, images_dir, width)
                elif child.name:
                    # Recursively process other elements
                    process_html_to_docx(str(child), doc, unique_id)

def process_element_content(element, para, images_dir):
    """Process content within an element with styles and images"""
    for content in element.contents:
        if content.name == 'img':
            # Handle inline images with width
            width = None
            if 'style' in content.attrs:
                style = content['style']
                width_match = re.search(r'width:\s*(\d+)px', style)
                if width_match:
                    width = int(width_match.group(1))
            add_image_to_paragraph(content, para, images_dir, width)
        elif content.name == 'span':
            run = para.add_run(content.text)
            apply_styles_to_run(run, content.get('style', ''))
        elif content.name in ['b', 'strong']:
            run = para.add_run(content.text)
            run.bold = True
        elif content.name in ['i', 'em']:
            run = para.add_run(content.text)
            run.italic = True
        elif content.name == 'u':
            run = para.add_run(content.text)
            run.underline = True
        elif content.name == 'br':
            para.add_run().add_break()
        elif content.name is None:  # Text node
            text = str(content).replace('\n', ' ')
            if text.strip():
                para.add_run(text)

from docx.shared import RGBColor

def apply_styles(run, style_str):
    """Apply CSS styles to a docx run"""
    styles = style_str.split(';')
    for style in styles:
        if not style.strip():
            continue
        parts = style.split(':', 1)
        if len(parts) != 2:
            continue

        prop, value = parts
        prop = prop.strip().lower()
        value = value.strip().lower()

        if prop == 'font-weight' and value == 'bold':
            run.bold = True
        elif prop == 'font-style' and value == 'italic':
            run.italic = True
        elif prop == 'text-decoration':
            if 'underline' in value:
                run.underline = True
        elif prop == 'color':
            match = re.search(r'#?([0-9a-f]{6})', value)
            if match:
                hex_color = match.group(1).upper()
                run.font.color.rgb = RGBColor.from_string(hex_color)


def process_image(element, parent, images_dir):
    """Process image element"""
    img_data = element['src'].split('base64,')[-1]
    img_bytes = base64.b64decode(img_data)
    img_filename = f"{uuid.uuid4()}.png"
    img_path = os.path.join(images_dir, img_filename)
    
    with open(img_path, 'wb') as f:
        f.write(img_bytes)
    
    if isinstance(parent, Document):
        parent.add_picture(img_path)
    else:  # Paragraph
        parent.add_run().add_picture(img_path)

@app.route('/upload-image', methods=['POST'])
def upload_image():
    if 'image' not in request.files:
        return jsonify({'error': 'No image uploaded'}), 400
    
    file = request.files['image']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    try:
        filename = secure_filename(file.filename)
        unique_id = str(uuid.uuid4())
        image_filename = f"{unique_id}_{filename}"
        image_path = os.path.join(app.config['UPLOAD_FOLDER'], image_filename)
        file.save(image_path)
        
        # Return URL with resized marker
        with open(image_path, 'rb') as img_file:
            img_data = base64.b64encode(img_file.read()).decode('utf-8')
            img_src = f"data:image/{filename.split('.')[-1]};base64,{img_data}"
            return jsonify({
                'url': img_src,
                'resized': True
            })
            
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/download-pdf-file/<filename>')
def download_pdf_file(filename):
    try:
        filename = secure_filename(filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        if not os.path.exists(file_path):
            return jsonify({'error': 'File not found'}), 404
            
        return send_from_directory(
            app.config['UPLOAD_FOLDER'],
            filename,
            as_attachment=True,
            mimetype='application/pdf'
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

        
# Background thread: Auto-cleanup old files
def cleanup_old_files(folder_paths, max_age_minutes=2880, check_interval_seconds=86400):
    """
    Aggressively clean up ALL files and folders in specified paths
    
    Args:
        folder_paths: List of folders to clean (e.g., ['uploads', 'converted'])
        max_age_minutes: Max age in minutes before deletion (default 120 = 2 hours)
        check_interval_seconds: How often to run cleanup (default 7200 = 2 hours)
    """
    def run_cleanup():
        while True:
            now = time.time()
            # Get the directory where app.py is located
            current_dir = os.path.dirname(os.path.abspath(__file__))
            
            # Process both the original paths and the paths relative to app.py
            all_paths = []
            for folder in folder_paths:
                # Add the original path (parent directory)
                all_paths.append(folder)
                # Add the path relative to app.py's directory
                all_paths.append(os.path.join(current_dir, folder))
            
            for folder in all_paths:
                try:
                    # Skip if folder doesn't exist
                    if not os.path.exists(folder):
                        continue
                        
                    # Process all items in folder
                    for item in os.listdir(folder):
                        item_path = os.path.join(folder, item)
                        
                        # Calculate age in minutes
                        age_minutes = (now - os.path.getmtime(item_path)) / 60
                        
                        if age_minutes > max_age_minutes:
                            try:
                                if os.path.isfile(item_path) or os.path.islink(item_path):
                                    os.remove(item_path)
                                    print(f"Deleted file: {item_path}")
                                elif os.path.isdir(item_path):
                                    shutil.rmtree(item_path)
                                    print(f"Deleted folder: {item_path}")
                            except Exception as e:
                                print(f"Error deleting {item_path}: {e}")
                except Exception as e:
                    print(f"Error processing folder {folder}: {e}")
                    
            time.sleep(check_interval_seconds)

    thread = threading.Thread(target=run_cleanup, daemon=True)
    thread.start()


    
# Default route to display welcome page
@app.route('/')
def welcome():
    return """
    <!DOCTYPE html>
    <html>
    <head>
        <title>PDFSimba Backend</title>
        <style>
            body {
                font-family: Arial, sans-serif;
                line-height: 1.6;
                color: #333;
                max-width: 800px;
                margin: 0 auto;
                padding: 20px;
                text-align: center;
            }
            h1 {
                color: #2c3e50;
                font-size: 2.5em;
                margin-bottom: 10px;
            }
            .logo {
                margin-bottom: 20px;
                font-size: 5em;
            }
            .status {
                background-color: #f1f8e9;
                border-left: 5px solid #7cb342;
                padding: 15px;
                margin-bottom: 20px;
                text-align: left;
            }
            .endpoints {
                background-color: #e8f4fd;
                border-left: 5px solid #2196F3;
                padding: 15px;
                text-align: left;
            }
            .endpoint {
                margin-bottom: 10px;
                font-family: monospace;
            }
        </style>
    </head>
    <body>
        <div class="logo">📄✨</div>
        <h1>Welcome to PDFSimba Backend</h1>
        <p>The PDF conversion service is running successfully.</p>
        
        <div class="status">
            <h2>Server Status</h2>
            <p>The backend server is online and ready to process requests.</p>
        </div>
        
        <div class="endpoints">
            <h2>Available Endpoints</h2>
            <div class="endpoint">/convert-word-to-pdf</div>
            <div class="endpoint">/convert-excel-to-pdf</div>
            <div class="endpoint">/convert-ppt-to-pdf</div>
            <div class="endpoint">/convert-jpg-to-pdf</div>
            <div class="endpoint">/convert-pdf-to-word</div>
            <div class="endpoint">/convert-pdf-to-excel</div>
            <div class="endpoint">/convert-pdf-to-ppt</div>
            <div class="endpoint">/convert-pdf-to-jpg</div>
            <div class="endpoint">/convert-pdf-to-png</div>
            <p>And many more...</p>
        </div>
        
        <p>This is a backend service. Please use the frontend application to interact with these services.</p>
    </body>
    </html>
    """


# Run server
if __name__ == '__main__':
    cleanup_old_files([UPLOAD_FOLDER, OUTPUT_FOLDER])
    app.run(debug=True, host='0.0.0.0', port=5000)
