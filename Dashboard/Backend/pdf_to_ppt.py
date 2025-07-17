import os
import uuid
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_FOLDER = os.path.abspath(os.path.join(BASE_DIR, '..', 'uploads'))
OUTPUT_FOLDER = os.path.abspath(os.path.join(BASE_DIR, '..', 'converted'))

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

def convert_pdf_to_ppt(file):
    if not file:
        return None, "No file uploaded."

    ext = os.path.splitext(file.filename)[1].lower()
    if ext != '.pdf':
        return None, "Only PDF files are supported."

    filename = f"{uuid.uuid4()}.pdf"
    input_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(input_path)

    output_filename = f"{os.path.splitext(filename)[0]}.pptx"
    output_path = os.path.join(OUTPUT_FOLDER, output_filename)

    try:
        images = convert_from_path(input_path)
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]  # Blank

        for image in images:
            slide = prs.slides.add_slide(blank_slide_layout)
            image_path = os.path.join(UPLOAD_FOLDER, f"temp_{uuid.uuid4()}.jpg")
            image.save(image_path, "JPEG")
            slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=prs.slide_width)
            os.remove(image_path)

        prs.save(output_path)
        return output_filename, None
    except Exception as e:
        return None, str(e)
